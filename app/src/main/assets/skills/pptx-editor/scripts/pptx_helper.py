"""Pptx-editor helper - programmatic ops on .pptx files.

Subcommands (all stateless, deterministic):
  inspect <file.pptx>                              - Print a per-slide shape
                                                     table (slide/shape-idx/role/
                                                     type/bbox/text) so the LLM
                                                     can see template layout.
  extract <file.pptx>                              - Plain text dump
                                                     (markitdown-style).
  set-text <file> --slide N --shape K --text "..." - Replace the text of one
      [--out OUT]                                    shape, PRESERVING the
                                                     original run-level format.
                                                     Default out: overwrite in
                                                     place. K is the shape-idx
                                                     from `inspect`.
  set-cell <file> --slide N --shape K --row R      - Same as set-text but for a
      --col C --text "..." [--out OUT]               table cell.
  new-from-outline <outline.json> --out OUT.pptx   - Build a 16:9 deck from an
      [--palette NAME] [--style NAME]                outline JSON using built-in
                                                     5 page-type recipes.

Design principles:
  1. Never assume `slide.shapes.title` or `placeholders[1]` exists. Always ask
     the LLM to `inspect` first and address shapes by (slide_idx, shape_idx).
  2. Preserve original run-level formatting on set-text / set-cell. This keeps
     template fidelity (font family, size, bold, color) instead of blowing it
     away like `shape.text = "..."` does.
  3. Built-in palettes + style recipes + page-type layouts let the LLM get a
     design-consistent deck with a JSON outline only.

Dependency: python-pptx (only).
"""
import sys
import os
import io
import json
import argparse
from copy import deepcopy

# Force UTF-8 stdout/stderr so emoji/CJK work in Windows/Android hosts.
try:
    if sys.stdout.encoding and sys.stdout.encoding.lower() not in ("utf-8", "utf8"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    if sys.stderr.encoding and sys.stderr.encoding.lower() not in ("utf-8", "utf8"):
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ===========================================================================
# Palettes (6-char hex -> RGBColor) - ONLY a handful; LLM can use custom too.
# ===========================================================================
PALETTES = {
    "business": {
        "primary":   "#2B2D42", "secondary": "#8D99AE", "accent": "#EF233C",
        "light":     "#EDF2F4", "bg":        "#FFFFFF",
    },
    "tech": {
        "primary":   "#03045E", "secondary": "#0077B6", "accent": "#FFB703",
        "light":     "#CAF0F8", "bg":        "#FFFFFF",
    },
    "wellness": {
        "primary":   "#006D77", "secondary": "#83C5BE", "accent": "#E29578",
        "light":     "#EDF6F9", "bg":        "#FFDDD2",
    },
    "elegant": {
        "primary":   "#22223B", "secondary": "#4A4E69", "accent": "#9A8C98",
        "light":     "#C9ADA7", "bg":        "#F2E9E4",
    },
    "education": {
        "primary":   "#264653", "secondary": "#2A9D8F", "accent": "#E76F51",
        "light":     "#E9C46A", "bg":        "#FFFFFF",
    },
    "platinum": {
        "primary":   "#0A0A0A", "secondary": "#525252", "accent": "#D4AF37",
        "light":     "#F5F5F5", "bg":        "#FFFFFF",
    },
}

# Four style recipes driving corner radius and spacing (inches).
STYLES = {
    "sharp":   {"radius": 0.02, "margin": 0.3, "gap": 0.2, "pad": 0.15},
    "soft":    {"radius": 0.08, "margin": 0.4, "gap": 0.3, "pad": 0.2},
    "rounded": {"radius": 0.18, "margin": 0.5, "gap": 0.4, "pad": 0.3},
    "pill":    {"radius": 0.30, "margin": 0.6, "gap": 0.5, "pad": 0.35},
}

SLIDE_W_IN = 13.333   # 16:9 width in inches (matches prs default 13.33x7.5)
SLIDE_H_IN = 7.5


# ===========================================================================
# Hard limits for the direct add-* CLI surface. The model fills strings; the
# helper auto-truncates / pads / clamps any field that exceeds these. Anything
# outside [min, max] becomes a [WARN] line in stdout, NEVER a failure.
# These numbers come from the actual rendered slide geometry (font size *
# available width / height) — beyond them the slide overflows visually.
# ===========================================================================
_LIM = {
    # Strings (max chars; longer -> trim with ellipsis)
    "title": 40,
    "subtitle": 80,
    "meta": 60,
    "section_intro": 120,
    "section_number": 4,
    "bullet": 80,
    "toc_item_title": 40,
    "toc_item_desc": 60,
    "stat_value": 8,
    "stat_label": 20,
    "stat_desc": 40,
    "step_label": 16,
    "step_desc": 60,
    "compare_col_title": 30,
    "compare_point": 60,
    "takeaway": 80,
    "contact": 60,
    "caption": 80,
    "header_cell": 16,
    "table_cell": 24,
    # Lists (count; longer -> drop tail; shorter -> let user provide)
    "bullets_min": 1, "bullets_max": 5,
    "headers_max": 5, "rows_max": 8,
    "stats_min": 2, "stats_max": 4,
    "steps_min": 3, "steps_max": 5,
    "takeaways_min": 2, "takeaways_max": 5,
    "toc_items_min": 3, "toc_items_max": 6,
    "compare_points_min": 2, "compare_points_max": 5,
}


def _clip_str(s, max_chars, *, where=None):
    """Truncate string to max_chars; print [WARN] if trimmed."""
    if s is None:
        return ""
    s = str(s)
    if len(s) <= max_chars:
        return s
    if where:
        print(f"[WARN] {where}: truncated from {len(s)} to {max_chars} chars")
    return s[: max_chars - 1] + "\u2026"


def _clip_list(lst, max_n, *, where=None):
    """Drop tail items beyond max_n; print [WARN] if dropped."""
    if not lst:
        return []
    out = list(lst)
    if len(out) <= max_n:
        return out
    if where:
        print(f"[WARN] {where}: kept first {max_n} of {len(out)} items")
    return out[:max_n]


def _split_csv(s, sep=","):
    """Split a CSV-like string with backslash-escape support: '\\,' -> ','.

    Returns trimmed non-empty parts. Tolerates None / list (passes through).
    """
    if s is None:
        return []
    if isinstance(s, list):
        return [str(x).strip() for x in s if str(x).strip()]
    s = str(s)
    parts = []
    cur = []
    i = 0
    while i < len(s):
        c = s[i]
        if c == "\\" and i + 1 < len(s) and s[i + 1] == sep:
            cur.append(sep)
            i += 2
            continue
        if c == sep:
            parts.append("".join(cur).strip())
            cur = []
            i += 1
            continue
        cur.append(c)
        i += 1
    parts.append("".join(cur).strip())
    return [p for p in parts if p != ""]


def _split_kv(s, sep=";", n_fields=2, defaults=("", "")):
    """Split 'a;b' or 'a;b;c' into a fixed-length tuple, padding with defaults.

    Used for --stat "value;label" / --step "label;desc" / --item "title;desc".
    """
    parts = _split_csv(s, sep=sep)
    out = list(parts[:n_fields])
    while len(out) < n_fields:
        out.append(defaults[len(out)] if len(out) < len(defaults) else "")
    return tuple(out)


# ===========================================================================
# Deck-level metadata: theme + page-num counter + section counter, encoded as
# JSON in core_properties.comments (with subject=_META_MARK as a sentinel).
# Survives normal save/reopen cycles — no sidecar file needed.
# ===========================================================================
_META_MARK = "pptx-helper-deck"


def _meta_default():
    return {"theme": "business", "next_page_num": 1, "section_count": 0}


def _meta_load(prs):
    cp = prs.core_properties
    raw = cp.comments or ""
    try:
        if cp.subject == _META_MARK and raw.startswith("{"):
            m = json.loads(raw)
            out = _meta_default()
            out.update({k: m[k] for k in out.keys() if k in m})
            # Sanitize types
            out["theme"] = str(out["theme"])
            out["next_page_num"] = int(out["next_page_num"])
            out["section_count"] = int(out["section_count"])
            return out
    except Exception:
        pass
    # No marker -> infer from existing slide count for forward compat
    fb = _meta_default()
    fb["next_page_num"] = max(1, len(prs.slides))
    return fb


def _meta_save(prs, meta):
    cp = prs.core_properties
    cp.subject = _META_MARK
    cp.comments = json.dumps(meta, ensure_ascii=False)


def _open_deck(path):
    if not os.path.exists(path):
        raise FileNotFoundError(
            "Deck not found: " + str(path) + "\n"
            "Hint: run `pptx_helper.py new-deck " + str(path) +
            " --theme business` first."
        )
    return Presentation(path)


# ===========================================================================
# Helpers
# ===========================================================================
def hex2rgb(h):
    """Accept '#FF0000' or 'FF0000' and return RGBColor."""
    h = h.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _emu_to_in(v):
    if v is None:
        return None
    return round(v / 914400.0, 2)


def _shape_role(shape):
    """Best-effort role label for LLM (title/body/subtitle/textbox/picture/...)."""
    if shape.is_placeholder:
        try:
            pf = shape.placeholder_format
            t = pf.type  # PP_PLACEHOLDER enum
            # Types: TITLE(13)=TITLE, BODY(2), CENTER_TITLE(15), SUBTITLE(4),
            #        OBJECT(7), PICTURE(18), CHART(12), TABLE(19), etc.
            name = str(t).split(".")[-1].lower() if t is not None else "placeholder"
            return f"ph.{name}(idx={pf.idx})"
        except Exception:
            return "placeholder"
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "textbox"
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "autoshape"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    return str(st).split(".")[-1].lower() if st is not None else "shape"


def _shape_text(shape, max_len=80):
    """Extract visible text (short preview) for inspect output."""
    if not shape.has_text_frame:
        # Table: concat first row + first col for preview
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                tbl = shape.table
                cells = []
                for r in range(min(tbl.rows.__len__(), 2)):
                    row = [tbl.cell(r, c).text.strip() for c in range(tbl.columns.__len__())]
                    cells.append(" | ".join(row))
                s = " // ".join(cells)
            except Exception:
                s = ""
        else:
            s = ""
    else:
        s = shape.text_frame.text.strip()
    s = s.replace("\n", "\\n").replace("\r", "")
    if len(s) > max_len:
        s = s[:max_len - 1] + "…"
    return s


# ===========================================================================
# Subcommand: inspect
# ===========================================================================
def cmd_inspect(args):
    prs = Presentation(args.file)
    sw = _emu_to_in(prs.slide_width)
    sh = _emu_to_in(prs.slide_height)
    print(f"[INSPECT] {args.file}")
    print(f"  slide_size: {sw} x {sh} in  ({len(prs.slides)} slides)")
    print()
    for sidx, slide in enumerate(prs.slides):
        layout = slide.slide_layout.name if slide.slide_layout else "?"
        print(f"─── Slide {sidx}  (layout: {layout}) ───")
        header = f"  {'idx':>3}  {'role':<22}  {'xy (in)':<14}  {'wh (in)':<14}  text"
        print(header)
        print("  " + "-" * (len(header) - 2))
        for shidx, shape in enumerate(slide.shapes):
            role = _shape_role(shape)
            x = _emu_to_in(shape.left)
            y = _emu_to_in(shape.top)
            w = _emu_to_in(shape.width)
            h = _emu_to_in(shape.height)
            xy = f"({x},{y})" if x is not None else "(-,-)"
            wh = f"{w}x{h}" if w is not None else "-"
            text = _shape_text(shape)
            print(f"  {shidx:>3}  {role:<22}  {xy:<14}  {wh:<14}  {text}")
            # If this shape is a table, enumerate cells so LLM can use set-cell
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                try:
                    tbl = shape.table
                    nr, nc = len(tbl.rows), len(tbl.columns)
                    print(f"        table: {nr} rows x {nc} cols")
                    for r in range(nr):
                        for c in range(nc):
                            t = tbl.cell(r, c).text.strip().replace("\n", "\\n")
                            if len(t) > 40:
                                t = t[:39] + "…"
                            print(f"          ({r},{c}): {t}")
                except Exception as e:
                    print(f"        table_err: {e}")
        print()
    _next_step([
        "You now have (slide_idx, shape_idx, role, text) for every shape.",
        f"Replace text: python pptx_helper.py set-text {args.file} --slide N --shape K --text \"...\"",
        f"Replace cell: python pptx_helper.py set-cell {args.file} --slide N --shape K --row R --col C --text \"...\"",
    ])


# ===========================================================================
# Subcommand: extract
# ===========================================================================
def _iter_text(shape):
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            line = "".join(r.text for r in p.runs) or p.text
            if line.strip():
                yield line
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        tbl = shape.table
        for r in range(len(tbl.rows)):
            row = [tbl.cell(r, c).text.strip() for c in range(len(tbl.columns))]
            yield " | ".join(row)
    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            yield from _iter_text(sub)


def cmd_extract(args):
    prs = Presentation(args.file)
    for sidx, slide in enumerate(prs.slides):
        print(f"### Slide {sidx + 1}")
        for shape in slide.shapes:
            for line in _iter_text(shape):
                print(line)
        print()
    _next_step([
        "Plain text dump above; use for summarization or search.",
        f"To locate which shape holds which text: python pptx_helper.py inspect {args.file}",
    ])


# ===========================================================================
# Subcommand: set-text  (preserves original run-level formatting)
# ===========================================================================
def _get_shape_by_idx(slide, shape_idx):
    shapes = list(slide.shapes)
    if shape_idx < 0 or shape_idx >= len(shapes):
        raise IndexError(f"shape idx {shape_idx} out of range (slide has {len(shapes)} shapes)")
    return shapes[shape_idx]


def _preserve_set_text_frame(tf, new_text):
    """Replace tf's text while preserving the first run's font formatting.

    Strategy:
      - split new_text by '\n' into lines
      - keep paragraph[0] as the format reference (clone its first run's font)
      - remove all other paragraphs (XML level) and rebuild
    """
    if not new_text:
        new_text = ""
    lines = new_text.split("\n")

    if len(tf.paragraphs) == 0:
        tf.text = new_text
        return

    p0 = tf.paragraphs[0]
    # Snapshot the first run's rPr XML (if any) so we can clone it.
    ref_rPr_xml = None
    if p0.runs:
        r0 = p0.runs[0]
        rPr_el = r0._r.find("{http://schemas.openxmlformats.org/drawingml/2006/main}rPr")
        if rPr_el is not None:
            ref_rPr_xml = deepcopy(rPr_el)

    # Snapshot the paragraph's pPr (alignment, indent, bullet) to reuse.
    pPr_el = p0._pPr
    ref_pPr_xml = deepcopy(pPr_el) if pPr_el is not None else None

    # Remove every paragraph except the first.
    txBody = tf._txBody
    for p in list(txBody):
        # Keep <a:bodyPr>/<a:lstStyle>; only remove <a:p>
        if p.tag.endswith("}p"):
            txBody.remove(p)

    # Re-add one <a:p> per line with cloned formatting.
    from pptx.oxml.ns import qn
    for i, line in enumerate(lines):
        p_el = txBody.makeelement(qn("a:p"), {})
        if ref_pPr_xml is not None:
            p_el.append(deepcopy(ref_pPr_xml))
        r_el = p_el.makeelement(qn("a:r"), {})
        if ref_rPr_xml is not None:
            r_el.append(deepcopy(ref_rPr_xml))
        t_el = r_el.makeelement(qn("a:t"), {})
        t_el.text = line
        r_el.append(t_el)
        p_el.append(r_el)
        txBody.append(p_el)


def cmd_set_text(args):
    # Compat shim: positional slide_pos/shape_pos/text_pos -> legacy attrs.
    def _pi(name, raw):
        try:
            return int(raw)
        except (TypeError, ValueError):
            raise ValueError(
                f"set-text: {name} must be an integer (got {raw!r}). Usage:\n"
                "  set-text FILE SLIDE SHAPE TEXT"
            )
    if args.slide is None and getattr(args, "slide_pos", None) is not None:
        args.slide = _pi("SLIDE", args.slide_pos)
    if args.shape is None and getattr(args, "shape_pos", None) is not None:
        args.shape = _pi("SHAPE", args.shape_pos)
    if not args.text and getattr(args, "text_pos", None) is not None:
        args.text = args.text_pos
    if args.slide is None or args.shape is None or args.text is None:
        raise ValueError(
            "set-text: SLIDE, SHAPE, TEXT all required. Usage:\n"
            "  set-text FILE SLIDE SHAPE TEXT"
        )
    prs = Presentation(args.file)
    if args.slide < 0 or args.slide >= len(prs.slides):
        raise IndexError(f"slide {args.slide} out of range (0..{len(prs.slides)-1})")
    slide = prs.slides[args.slide]
    shape = _get_shape_by_idx(slide, args.shape)
    if not shape.has_text_frame:
        raise TypeError(f"shape {args.shape} on slide {args.slide} has no text_frame "
                        f"(role={_shape_role(shape)}). Use set-cell for tables.")
    _preserve_set_text_frame(shape.text_frame, args.text)
    out = args.out or args.file
    prs.save(out)
    print(f"[OK] set-text slide={args.slide} shape={args.shape} -> {out}")
    _next_step([
        f"Text updated on slide {args.slide} shape {args.shape}.",
        f"Verify:   python pptx_helper.py inspect {out}",
        "Continue editing other shapes with more set-text calls.",
    ])


# ===========================================================================
# Subcommand: set-cell  (table cell preserve-format replace)
# ===========================================================================
def cmd_set_cell(args):
    # Compat shim: positional slide_pos/shape_pos/row_pos/col_pos/text_pos.
    def _pi(name, raw):
        try:
            return int(raw)
        except (TypeError, ValueError):
            raise ValueError(
                f"set-cell: {name} must be an integer (got {raw!r}). Usage:\n"
                "  set-cell FILE SLIDE SHAPE ROW COL TEXT"
            )
    if args.slide is None and getattr(args, "slide_pos", None) is not None:
        args.slide = _pi("SLIDE", args.slide_pos)
    if args.shape is None and getattr(args, "shape_pos", None) is not None:
        args.shape = _pi("SHAPE", args.shape_pos)
    if args.row is None and getattr(args, "row_pos", None) is not None:
        args.row = _pi("ROW", args.row_pos)
    if args.col is None and getattr(args, "col_pos", None) is not None:
        args.col = _pi("COL", args.col_pos)
    if not args.text and getattr(args, "text_pos", None) is not None:
        args.text = args.text_pos
    if (args.slide is None or args.shape is None or args.row is None
            or args.col is None or args.text is None):
        raise ValueError(
            "set-cell: SLIDE, SHAPE, ROW, COL, TEXT all required. Usage:\n"
            "  set-cell FILE SLIDE SHAPE ROW COL TEXT"
        )
    prs = Presentation(args.file)
    if args.slide < 0 or args.slide >= len(prs.slides):
        raise IndexError(f"slide {args.slide} out of range")
    slide = prs.slides[args.slide]
    shape = _get_shape_by_idx(slide, args.shape)
    if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        raise TypeError(f"shape {args.shape} on slide {args.slide} is not a table "
                        f"(role={_shape_role(shape)})")
    tbl = shape.table
    if args.row < 0 or args.row >= len(tbl.rows):
        raise IndexError(f"row {args.row} out of range (0..{len(tbl.rows)-1})")
    if args.col < 0 or args.col >= len(tbl.columns):
        raise IndexError(f"col {args.col} out of range (0..{len(tbl.columns)-1})")
    cell = tbl.cell(args.row, args.col)
    _preserve_set_text_frame(cell.text_frame, args.text)
    out = args.out or args.file
    prs.save(out)
    print(f"[OK] set-cell slide={args.slide} shape={args.shape} "
          f"cell=({args.row},{args.col}) -> {out}")
    _next_step([
        f"Cell ({args.row},{args.col}) updated on slide {args.slide} shape {args.shape}.",
        f"Verify:   python pptx_helper.py inspect {out}",
    ])


# ===========================================================================
# Subcommand: new-from-outline  (5 page-type deck builder)
# ===========================================================================
def _solid_fill(shape, rgb_hex):
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex2rgb(rgb_hex)


def _no_line(shape):
    try:
        shape.line.fill.background()
    except Exception:
        try:
            shape.line.width = 0
        except Exception:
            pass


def _add_textbox(slide, x, y, w, h, text, *, font_size=14, bold=False,
                 color_hex="#0A0A0A", align=None, valign=None, font_name=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if valign is not None:
        tf.vertical_anchor = valign
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if align is not None:
            p.alignment = align
        # After setting p.text a run exists.
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.bold = bold
            r.font.color.rgb = hex2rgb(color_hex)
            if font_name:
                r.font.name = font_name
    return tb


def _add_page_badge(slide, n, accent_hex, style_pad=0.0):
    """Bottom-right page number pill badge. Skip on cover."""
    x = SLIDE_W_IN - 0.7
    y = SLIDE_H_IN - 0.5
    w, h = 0.4, 0.3
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    _solid_fill(pill, accent_hex)
    _no_line(pill)
    tf = pill.text_frame
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.0);   tf.margin_bottom = Inches(0.0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(n)
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = hex2rgb("#FFFFFF")


def _add_bg(slide, hex_color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    _solid_fill(bg, hex_color)
    _no_line(bg)
    return bg


# --- Page-type builders ----------------------------------------------------
def _build_cover(slide, data, pal, style):
    _add_bg(slide, pal["primary"])
    # Big title
    _add_textbox(slide, 1.0, 2.6, SLIDE_W_IN - 2.0, 1.6,
                 data.get("title", "Title"),
                 font_size=54, bold=True, color_hex="#FFFFFF",
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if data.get("subtitle"):
        _add_textbox(slide, 1.0, 4.3, SLIDE_W_IN - 2.0, 0.7,
                     data["subtitle"], font_size=22, bold=False,
                     color_hex=pal["light"], align=PP_ALIGN.CENTER)
    if data.get("meta"):
        _add_textbox(slide, 1.0, 6.6, SLIDE_W_IN - 2.0, 0.4,
                     data["meta"], font_size=14, bold=False,
                     color_hex=pal["light"], align=PP_ALIGN.CENTER)


def _build_toc(slide, data, pal, style, page_num):
    _add_bg(slide, pal["bg"])
    _add_textbox(slide, style["margin"], 0.5, SLIDE_W_IN - 2 * style["margin"], 0.8,
                 data.get("title", "Table of Contents"),
                 font_size=36, bold=True, color_hex=pal["primary"])
    sections = data.get("sections", [])
    n = len(sections)
    top = 1.8
    row_h = min(0.9, (SLIDE_H_IN - top - 1.2) / max(n, 1))
    for i, sec in enumerate(sections):
        y = top + i * row_h
        num = f"{i+1:02d}"
        # Number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(style["margin"]), Inches(y),
                                       Inches(0.5), Inches(0.5))
        _solid_fill(badge, pal["accent"])
        _no_line(badge)
        tf = badge.text_frame
        tf.margin_left = Inches(0); tf.margin_right = Inches(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(12); r.font.bold = True
            r.font.color.rgb = hex2rgb("#FFFFFF")
        # Section title + desc
        if isinstance(sec, str):
            title, desc = sec, ""
        else:
            title = sec.get("title", ""); desc = sec.get("desc", "")
        _add_textbox(slide, style["margin"] + 0.8, y - 0.02,
                     SLIDE_W_IN - 2 * style["margin"] - 0.8, 0.45,
                     title, font_size=22, bold=True, color_hex=pal["primary"])
        if desc:
            _add_textbox(slide, style["margin"] + 0.8, y + 0.4,
                         SLIDE_W_IN - 2 * style["margin"] - 0.8, 0.35,
                         desc, font_size=13, color_hex=pal["secondary"])
    _add_page_badge(slide, page_num, pal["accent"])


def _build_section(slide, data, pal, style, page_num):
    _add_bg(slide, pal["primary"])
    # Accent block on the left
    blk = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(1.5), Inches(SLIDE_H_IN))
    _solid_fill(blk, pal["accent"])
    _no_line(blk)
    num = data.get("number", "01")
    _add_textbox(slide, 2.0, 2.3, SLIDE_W_IN - 2.3, 1.5, num,
                 font_size=96, bold=True, color_hex=pal["accent"],
                 valign=MSO_ANCHOR.MIDDLE)
    _add_textbox(slide, 2.0, 3.7, SLIDE_W_IN - 2.3, 1.2,
                 data.get("title", "Section Title"),
                 font_size=42, bold=True, color_hex="#FFFFFF")
    if data.get("intro"):
        _add_textbox(slide, 2.0, 5.0, SLIDE_W_IN - 2.3, 1.0,
                     data["intro"], font_size=16, color_hex=pal["light"])
    _add_page_badge(slide, page_num, pal["accent"])


def _build_content_text(slide, data, pal, style, page_num):
    _add_bg(slide, pal["bg"])
    _add_textbox(slide, style["margin"], 0.5, SLIDE_W_IN - 2 * style["margin"], 0.8,
                 data.get("title", ""), font_size=32, bold=True,
                 color_hex=pal["primary"])
    bullets = data.get("bullets", [])
    top = 1.8
    w = SLIDE_W_IN - 2 * style["margin"]
    line_h = 0.55
    for i, b in enumerate(bullets[:8]):
        y = top + i * line_h
        # Bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(style["margin"]), Inches(y + 0.18),
                                     Inches(0.14), Inches(0.14))
        _solid_fill(dot, pal["accent"])
        _no_line(dot)
        _add_textbox(slide, style["margin"] + 0.3, y, w - 0.3, 0.5,
                     b, font_size=18, color_hex=pal["primary"])
    _add_page_badge(slide, page_num, pal["accent"])


def _build_content_bullets_image(slide, data, pal, style, page_num):
    _add_bg(slide, pal["bg"])
    _add_textbox(slide, style["margin"], 0.5, SLIDE_W_IN - 2 * style["margin"], 0.8,
                 data.get("title", ""), font_size=32, bold=True,
                 color_hex=pal["primary"])
    # Left column bullets
    bullets = data.get("bullets", [])
    col_w = (SLIDE_W_IN - 2 * style["margin"] - style["gap"]) / 2
    top = 1.8
    line_h = 0.55
    for i, b in enumerate(bullets[:6]):
        y = top + i * line_h
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(style["margin"]), Inches(y + 0.18),
                                     Inches(0.14), Inches(0.14))
        _solid_fill(dot, pal["accent"])
        _no_line(dot)
        _add_textbox(slide, style["margin"] + 0.3, y, col_w - 0.3, 0.5,
                     b, font_size=16, color_hex=pal["primary"])
    # Right column image (if exists)
    img = data.get("image")
    img_x = style["margin"] + col_w + style["gap"]
    img_y = top
    img_w = col_w
    img_h = SLIDE_H_IN - img_y - 1.0
    if img and os.path.exists(img):
        try:
            slide.shapes.add_picture(img, Inches(img_x), Inches(img_y),
                                     width=Inches(img_w), height=Inches(img_h))
        except Exception as e:
            _add_textbox(slide, img_x, img_y, img_w, img_h,
                         f"[image err: {e}]", font_size=12,
                         color_hex=pal["secondary"])
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(img_x), Inches(img_y),
                                    Inches(img_w), Inches(img_h))
        ph.fill.solid(); ph.fill.fore_color.rgb = hex2rgb(pal["light"])
        _no_line(ph)
        _add_textbox(slide, img_x, img_y + img_h / 2 - 0.2, img_w, 0.4,
                     data.get("image_caption", "[Image Placeholder]"),
                     font_size=14, color_hex=pal["secondary"],
                     align=PP_ALIGN.CENTER)
    _add_page_badge(slide, page_num, pal["accent"])


def _build_content_table(slide, data, pal, style, page_num):
    _add_bg(slide, pal["bg"])
    _add_textbox(slide, style["margin"], 0.5, SLIDE_W_IN - 2 * style["margin"], 0.8,
                 data.get("title", ""), font_size=32, bold=True,
                 color_hex=pal["primary"])
    tdata = data.get("table", {})
    headers = tdata.get("headers", [])
    rows = tdata.get("rows", [])
    nc = len(headers)
    nr = len(rows) + 1
    if nc == 0 or nr == 1:
        return
    t_x = style["margin"]
    t_y = 1.6
    t_w = SLIDE_W_IN - 2 * style["margin"]
    t_h = min(SLIDE_H_IN - t_y - 1.0, 0.5 + 0.4 * (nr - 1))
    table = slide.shapes.add_table(nr, nc, Inches(t_x), Inches(t_y),
                                   Inches(t_w), Inches(t_h)).table
    # Header
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        cell.fill.solid(); cell.fill.fore_color.rgb = hex2rgb(pal["primary"])
        tf = cell.text_frame
        p = tf.paragraphs[0]; p.text = str(h)
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(14)
            r.font.color.rgb = hex2rgb("#FFFFFF")
    # Rows (zebra)
    for ri, row in enumerate(rows):
        bg = pal["light"] if ri % 2 == 0 else pal["bg"]
        for ci in range(nc):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = hex2rgb(bg)
            v = row[ci] if ci < len(row) else ""
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(v)
            for r in p.runs:
                r.font.size = Pt(12)
                r.font.color.rgb = hex2rgb(pal["primary"])
    _add_page_badge(slide, page_num, pal["accent"])


def _build_content_stats(slide, data, pal, style, page_num):
    _add_bg(slide, pal["bg"])
    _add_textbox(slide, style["margin"], 0.5, SLIDE_W_IN - 2 * style["margin"], 0.8,
                 data.get("title", ""), font_size=32, bold=True,
                 color_hex=pal["primary"])
    stats = data.get("stats", [])[:4]
    n = max(len(stats), 1)
    total_w = SLIDE_W_IN - 2 * style["margin"]
    col_w = (total_w - (n - 1) * style["gap"]) / n
    top = 2.2
    box_h = 3.0
    for i, s in enumerate(stats):
        x = style["margin"] + i * (col_w + style["gap"])
        # Big number
        _add_textbox(slide, x, top, col_w, 1.6,
                     s.get("number", ""),
                     font_size=66, bold=True, color_hex=pal["accent"],
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.BOTTOM)
        _add_textbox(slide, x, top + 1.7, col_w, 0.5,
                     s.get("label", ""),
                     font_size=16, bold=True, color_hex=pal["primary"],
                     align=PP_ALIGN.CENTER)
        if s.get("desc"):
            _add_textbox(slide, x, top + 2.3, col_w, 0.7,
                         s["desc"], font_size=12, color_hex=pal["secondary"],
                         align=PP_ALIGN.CENTER)
    _add_page_badge(slide, page_num, pal["accent"])


def _build_content_comparison(slide, data, pal, style, page_num):
    _add_bg(slide, pal["bg"])
    _add_textbox(slide, style["margin"], 0.5, SLIDE_W_IN - 2 * style["margin"], 0.8,
                 data.get("title", ""), font_size=32, bold=True,
                 color_hex=pal["primary"])
    left = data.get("left", {}); right = data.get("right", {})
    col_w = (SLIDE_W_IN - 2 * style["margin"] - style["gap"]) / 2
    top = 1.7
    col_h = SLIDE_H_IN - top - 1.0
    for i, col in enumerate([left, right]):
        cx = style["margin"] + i * (col_w + style["gap"])
        bg = pal["light"] if i == 0 else pal["secondary"]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(cx), Inches(top),
                                      Inches(col_w), Inches(col_h))
        _solid_fill(card, bg)
        _no_line(card)
        text_color = pal["primary"] if i == 0 else "#FFFFFF"
        _add_textbox(slide, cx + 0.3, top + 0.3, col_w - 0.6, 0.6,
                     col.get("title", ""), font_size=22, bold=True,
                     color_hex=text_color)
        for bi, b in enumerate(col.get("points", [])[:6]):
            _add_textbox(slide, cx + 0.3, top + 1.1 + bi * 0.5, col_w - 0.6, 0.45,
                         f"• {b}", font_size=14, color_hex=text_color)
    _add_page_badge(slide, page_num, pal["accent"])


def _build_summary(slide, data, pal, style, page_num):
    _add_bg(slide, pal["primary"])
    _add_textbox(slide, style["margin"], 0.8, SLIDE_W_IN - 2 * style["margin"], 1.0,
                 data.get("title", "Key Takeaways"),
                 font_size=44, bold=True, color_hex="#FFFFFF",
                 align=PP_ALIGN.CENTER)
    takeaways = data.get("takeaways", [])
    top = 2.4
    line_h = 0.55
    for i, t in enumerate(takeaways[:5]):
        y = top + i * line_h
        mark = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(style["margin"] + 1.5),
                                      Inches(y + 0.12), Inches(0.3), Inches(0.3))
        _solid_fill(mark, pal["accent"])
        _no_line(mark)
        check = mark.text_frame
        check.vertical_anchor = MSO_ANCHOR.MIDDLE
        pm = check.paragraphs[0]
        pm.text = "✓"; pm.alignment = PP_ALIGN.CENTER
        for r in pm.runs:
            r.font.size = Pt(14); r.font.bold = True
            r.font.color.rgb = hex2rgb("#FFFFFF")
        _add_textbox(slide, style["margin"] + 2.0, y,
                     SLIDE_W_IN - 2 * style["margin"] - 2.0, 0.5,
                     t, font_size=18, color_hex=pal["light"])
    if data.get("contact"):
        _add_textbox(slide, style["margin"], SLIDE_H_IN - 1.0,
                     SLIDE_W_IN - 2 * style["margin"], 0.4,
                     data["contact"], font_size=14,
                     color_hex=pal["light"], align=PP_ALIGN.CENTER)
    _add_page_badge(slide, page_num, pal["accent"])


# ---------------------------------------------------------------------------
# Plan-API extra builders (split text+table, split text+image, timeline).
# Used by cmd_plan via the LAYOUTS catalog. Reuse helpers above.
# ---------------------------------------------------------------------------
def _draw_table_at(slide, x, y, w, h, headers, rows, pal):
    """Helper: draw a styled table at (x, y, w, h) inches. No badge, no title."""
    nc = len(headers)
    nr = len(rows) + 1
    if nc == 0 or nr == 1:
        return
    actual_h = min(h, 0.5 + 0.4 * (nr - 1))
    table = slide.shapes.add_table(nr, nc, Inches(x), Inches(y),
                                   Inches(w), Inches(actual_h)).table
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        cell.fill.solid(); cell.fill.fore_color.rgb = hex2rgb(pal["primary"])
        p = cell.text_frame.paragraphs[0]; p.text = str(head)
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(13)
            r.font.color.rgb = hex2rgb("#FFFFFF")
    for ri, row in enumerate(rows):
        bg = pal["light"] if ri % 2 == 0 else pal["bg"]
        for ci in range(nc):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = hex2rgb(bg)
            v = row[ci] if ci < len(row) else ""
            cell.text = ""
            p = cell.text_frame.paragraphs[0]; p.text = str(v)
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.color.rgb = hex2rgb(pal["primary"])


def _build_text_table_split(slide, data, pal, style, page_num, *, table_pos):
    """Mixed text+table layout. table_pos in {'left','right','top','bottom'}.

    Slots (from data dict): title, bullets, headers, rows.
    """
    _add_bg(slide, pal["bg"])
    _add_textbox(slide, style["margin"], 0.5, SLIDE_W_IN - 2 * style["margin"], 0.8,
                 data.get("title", ""), font_size=32, bold=True,
                 color_hex=pal["primary"])
    margin = style["margin"]
    gap = style["gap"]
    bullets = data.get("bullets", [])
    headers = data.get("headers", [])
    rows = data.get("rows", [])
    if table_pos in ("left", "right"):
        # Side-by-side, 50/50 split
        col_w = (SLIDE_W_IN - 2 * margin - gap) / 2
        top = 1.7
        col_h = SLIDE_H_IN - top - 1.0
        if table_pos == "right":
            text_x, table_x = margin, margin + col_w + gap
        else:
            table_x, text_x = margin, margin + col_w + gap
        bs = bullets[:5]
        line_h = min(0.6, max(0.4, col_h / max(len(bs), 1)))
        for i, b in enumerate(bs):
            y = top + i * line_h
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(text_x), Inches(y + 0.18),
                                         Inches(0.14), Inches(0.14))
            _solid_fill(dot, pal["accent"]); _no_line(dot)
            _add_textbox(slide, text_x + 0.3, y, col_w - 0.3, line_h,
                         b, font_size=15, color_hex=pal["primary"])
        _draw_table_at(slide, table_x, top, col_w, col_h, headers, rows, pal)
    else:  # top / bottom (vertical stack)
        top = 1.7
        avail_h = SLIDE_H_IN - top - 0.8
        block_h = (avail_h - gap) / 2
        if table_pos == "bottom":
            text_y, table_y = top, top + block_h + gap
        else:
            table_y, text_y = top, top + block_h + gap
        full_w = SLIDE_W_IN - 2 * margin
        bs = bullets[:3]
        line_h = min(0.55, max(0.4, block_h / max(len(bs), 1)))
        for i, b in enumerate(bs):
            y = text_y + i * line_h
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(margin), Inches(y + 0.16),
                                         Inches(0.14), Inches(0.14))
            _solid_fill(dot, pal["accent"]); _no_line(dot)
            _add_textbox(slide, margin + 0.3, y, full_w - 0.3, line_h,
                         b, font_size=16, color_hex=pal["primary"])
        _draw_table_at(slide, margin, table_y, full_w, block_h,
                       headers, rows, pal)
    _add_page_badge(slide, page_num, pal["accent"])


def _build_text_image_split(slide, data, pal, style, page_num, *, image_pos):
    """Mixed text+image layout. image_pos in {'left','right'}.

    Slots: title, bullets, image_path (optional path), caption (optional).
    """
    _add_bg(slide, pal["bg"])
    _add_textbox(slide, style["margin"], 0.5, SLIDE_W_IN - 2 * style["margin"], 0.8,
                 data.get("title", ""), font_size=32, bold=True,
                 color_hex=pal["primary"])
    margin, gap = style["margin"], style["gap"]
    bullets = data.get("bullets", [])[:5]
    img_path = data.get("image_path") or data.get("image")
    col_w = (SLIDE_W_IN - 2 * margin - gap) / 2
    top = 1.7
    col_h = SLIDE_H_IN - top - 1.0
    if image_pos == "right":
        text_x, img_x = margin, margin + col_w + gap
    else:
        img_x, text_x = margin, margin + col_w + gap
    line_h = 0.55
    for i, b in enumerate(bullets):
        y = top + i * line_h
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(text_x), Inches(y + 0.18),
                                     Inches(0.14), Inches(0.14))
        _solid_fill(dot, pal["accent"]); _no_line(dot)
        _add_textbox(slide, text_x + 0.3, y, col_w - 0.3, 0.5,
                     b, font_size=15, color_hex=pal["primary"])
    if img_path and os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(img_x), Inches(top),
                                     width=Inches(col_w), height=Inches(col_h))
        except Exception as e:
            _add_textbox(slide, img_x, top, col_w, col_h,
                         f"[image err: {e}]", font_size=12,
                         color_hex=pal["secondary"])
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(img_x), Inches(top),
                                    Inches(col_w), Inches(col_h))
        ph.fill.solid(); ph.fill.fore_color.rgb = hex2rgb(pal["light"])
        _no_line(ph)
        _add_textbox(slide, img_x, top + col_h / 2 - 0.2, col_w, 0.4,
                     data.get("caption", "[Image Placeholder]"),
                     font_size=14, color_hex=pal["secondary"],
                     align=PP_ALIGN.CENTER)
    _add_page_badge(slide, page_num, pal["accent"])


def _build_content_timeline(slide, data, pal, style, page_num):
    """Horizontal 3-5 step timeline with numbered circles connected by a line.

    Slots: title, steps[] of {label, desc?}.
    """
    _add_bg(slide, pal["bg"])
    _add_textbox(slide, style["margin"], 0.5, SLIDE_W_IN - 2 * style["margin"], 0.8,
                 data.get("title", ""), font_size=32, bold=True,
                 color_hex=pal["primary"])
    steps = data.get("steps", [])[:5]
    n = max(len(steps), 1)
    margin = style["margin"]
    total_w = SLIDE_W_IN - 2 * margin
    col_w = total_w / n
    top = 2.5
    # Horizontal connector line behind the circles
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(margin + col_w / 2),
                                  Inches(top + 0.45),
                                  Inches(max(total_w - col_w, 0.1)),
                                  Inches(0.05))
    _solid_fill(line, pal["secondary"]); _no_line(line)
    for i, st in enumerate(steps):
        cx = margin + i * col_w
        cb = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(cx + col_w / 2 - 0.4),
                                    Inches(top), Inches(0.8), Inches(0.8))
        _solid_fill(cb, pal["accent"]); _no_line(cb)
        tf = cb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = str(i + 1); p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(22); r.font.bold = True
            r.font.color.rgb = hex2rgb("#FFFFFF")
        label = st.get("label", "") if isinstance(st, dict) else str(st)
        _add_textbox(slide, cx, top + 1.0, col_w, 0.5,
                     label, font_size=16, bold=True, color_hex=pal["primary"],
                     align=PP_ALIGN.CENTER)
        if isinstance(st, dict) and st.get("desc"):
            _add_textbox(slide, cx, top + 1.6, col_w, 1.5,
                         st["desc"], font_size=12, color_hex=pal["secondary"],
                         align=PP_ALIGN.CENTER)
    _add_page_badge(slide, page_num, pal["accent"])


_BUILDERS = {
    "cover":             _build_cover,
    "toc":               _build_toc,
    "section":           _build_section,
    "content.text":      _build_content_text,
    "content.bullets_image": _build_content_bullets_image,
    "content.table":     _build_content_table,
    "content.stats":     _build_content_stats,
    "content.comparison": _build_content_comparison,
    "summary":           _build_summary,
}

# Hard caps on content density (drives readable layout).
LIMITS = {
    "toc.sections":          (3, 6),    # min, max
    "content.text.bullets":  (2, 6),
    "content.bullets_image.bullets": (2, 6),
    "content.stats.stats":   (2, 4),
    "content.comparison.points": (2, 6),  # per side
    "content.table.rows":    (1, 8),
    "summary.takeaways":     (2, 5),
}


def _validate_outline(outline):
    """Return list of issue strings; empty list = OK. Soft warnings included."""
    issues = []
    if not isinstance(outline, dict):
        return ["outline must be a JSON object"]

    pal = outline.get("palette", "business")
    if isinstance(pal, str) and pal not in PALETTES:
        issues.append(f"palette '{pal}' unknown; choices: {list(PALETTES)}")
    style = outline.get("style", "soft")
    if style not in STYLES:
        issues.append(f"style '{style}' unknown; choices: {list(STYLES)}")

    slides = outline.get("slides", [])
    if not isinstance(slides, list):
        return issues + ["'slides' must be an array"]

    for i, sd in enumerate(slides):
        prefix = f"slides[{i}]"
        if not isinstance(sd, dict):
            issues.append(f"{prefix}: must be an object"); continue
        t = sd.get("type", "content")
        if t == "content":
            sub = sd.get("subtype", "text")
            key = f"content.{sub}"
        else:
            key = t
        if key not in _BUILDERS:
            issues.append(f"{prefix}: unknown type/subtype '{key}'")
            continue

        # Check density caps
        if key == "toc":
            n = len(sd.get("sections", []))
            lo, hi = LIMITS["toc.sections"]
            if not (lo <= n <= hi):
                issues.append(f"{prefix}: toc.sections={n}, recommend {lo}..{hi}")
        elif key == "content.text":
            n = len(sd.get("bullets", []))
            lo, hi = LIMITS["content.text.bullets"]
            if not (lo <= n <= hi):
                issues.append(f"{prefix}: content.text.bullets={n}, "
                              f"recommend {lo}..{hi} (will truncate at {hi})")
        elif key == "content.bullets_image":
            n = len(sd.get("bullets", []))
            lo, hi = LIMITS["content.bullets_image.bullets"]
            if not (lo <= n <= hi):
                issues.append(f"{prefix}: bullets_image.bullets={n}, "
                              f"recommend {lo}..{hi}")
            if sd.get("image") and not os.path.exists(sd["image"]):
                issues.append(f"{prefix}: image not found: {sd['image']} "
                              f"(will render placeholder)")
        elif key == "content.stats":
            n = len(sd.get("stats", []))
            lo, hi = LIMITS["content.stats.stats"]
            if not (lo <= n <= hi):
                issues.append(f"{prefix}: stats={n}, recommend {lo}..{hi} "
                              f"(will truncate at {hi})")
        elif key == "content.comparison":
            for side in ("left", "right"):
                pts = sd.get(side, {}).get("points", [])
                lo, hi = LIMITS["content.comparison.points"]
                if not (lo <= len(pts) <= hi):
                    issues.append(f"{prefix}: comparison.{side}.points="
                                  f"{len(pts)}, recommend {lo}..{hi}")
        elif key == "content.table":
            tdata = sd.get("table", {})
            nr = len(tdata.get("rows", []))
            lo, hi = LIMITS["content.table.rows"]
            if not (lo <= nr <= hi):
                issues.append(f"{prefix}: table.rows={nr}, recommend {lo}..{hi}")
            nh = len(tdata.get("headers", []))
            if nh == 0:
                issues.append(f"{prefix}: table.headers is empty")
        elif key == "summary":
            n = len(sd.get("takeaways", []))
            lo, hi = LIMITS["summary.takeaways"]
            if not (lo <= n <= hi):
                issues.append(f"{prefix}: takeaways={n}, recommend {lo}..{hi}")

    # Layout monotony check (3+ same subtype in a row)
    keys_seq = []
    for sd in slides:
        t = sd.get("type", "content")
        keys_seq.append(f"content.{sd.get('subtype', 'text')}" if t == "content" else t)
    run = 1
    for i in range(1, len(keys_seq)):
        if keys_seq[i] == keys_seq[i - 1] and keys_seq[i].startswith("content."):
            run += 1
            if run >= 3:
                issues.append(f"slides[{i-2}..{i}]: same subtype "
                              f"'{keys_seq[i]}' repeated 3+ times — vary layouts")
                run = 1
        else:
            run = 1

    return issues


# ===========================================================================
# Built-in example outlines (LLM can copy + edit fields, no need to invent
# structure). Use `examples list` / `examples <name>` to print them.
# ===========================================================================
EXAMPLES_PPTX = {
    "pitch": {
        "_desc": "8-page startup pitch deck (cover + problem + solution + market + traction + ask + summary).",
        "outline": {
            "title": "StartupX",
            "subtitle": "Revolutionizing offline AI",
            "meta": "Pitch Deck · 2025-Q1",
            "palette": "tech",
            "style": "rounded",
            "slides": [
                {"type": "section", "number": "01", "title": "Problem",
                 "intro": "用户在弱网/隐私敏感场景下无法用 AI"},
                {"type": "content", "subtype": "stats",
                 "title": "Market Pain",
                 "stats": [
                     {"number": "73%", "label": "Privacy Concern", "desc": "Pew 2024"},
                     {"number": "$50B", "label": "TAM by 2028", "desc": "Gartner"},
                     {"number": "2.5x", "label": "Latency Today", "desc": "vs offline"}]},
                {"type": "section", "number": "02", "title": "Our Solution"},
                {"type": "content", "subtype": "bullets_image",
                 "title": "On-Device LLM",
                 "bullets": ["100% offline inference",
                             "<200ms first token",
                             "Works on 4GB RAM phone",
                             "End-to-end encrypted by design"],
                 "image_caption": "[Architecture Diagram]"},
                {"type": "content", "subtype": "comparison",
                 "title": "vs. Cloud-Only",
                 "left":  {"title": "Cloud LLM",
                           "points": ["Latency 1-3s", "Privacy risk",
                                      "Per-call cost", "Network required"]},
                 "right": {"title": "Our Edge LLM",
                           "points": ["Latency <200ms", "Zero data leak",
                                      "Zero marginal cost", "Works offline"]}},
                {"type": "section", "number": "03", "title": "Traction"},
                {"type": "content", "subtype": "stats",
                 "title": "Last 6 Months",
                 "stats": [
                     {"number": "120K", "label": "MAU", "desc": "+18% MoM"},
                     {"number": "4.7", "label": "App Rating", "desc": "23K reviews"},
                     {"number": "$2M", "label": "ARR", "desc": "Q4 run rate"},
                     {"number": "8", "label": "Enterprise", "desc": "Logos"}]},
                {"type": "summary", "title": "The Ask",
                 "takeaways": ["$5M Series A",
                               "18-month runway to Series B",
                               "Hire 12 engineers in APAC",
                               "Launch enterprise tier in 2025-Q3"],
                 "contact": "founders@startupx.ai"}
            ]
        }
    },

    "quarterly": {
        "_desc": "10-page quarterly business review (highlights + financials + roadmap).",
        "outline": {
            "title": "Q4 2024 Business Review",
            "subtitle": "Performance & Roadmap",
            "meta": "Internal · 2024-12-31",
            "palette": "business",
            "style": "soft",
            "slides": [
                {"type": "toc", "title": "Agenda",
                 "sections": [
                     {"title": "Highlights", "desc": "Quarterly wins"},
                     {"title": "Financials", "desc": "Revenue and cost"},
                     {"title": "Products",  "desc": "Launches and iterations"},
                     {"title": "Roadmap",   "desc": "Next quarter focus"}]},
                {"type": "section", "number": "01", "title": "Highlights"},
                {"type": "content", "subtype": "stats",
                 "title": "By the Numbers",
                 "stats": [
                     {"number": "+18%", "label": "Revenue YoY"},
                     {"number": "62",   "label": "NPS",          "desc": "+17 pts"},
                     {"number": "3",    "label": "New Products"},
                     {"number": "58",   "label": "Team Size",    "desc": "+16 hires"}]},
                {"type": "content", "subtype": "text",
                 "title": "What Went Well",
                 "bullets": ["Revenue grew 18% YoY to $5.2M",
                             "Launched 3 major product lines on schedule",
                             "Expanded into APAC market",
                             "NPS improved from 45 to 62"]},
                {"type": "section", "number": "02", "title": "Financials"},
                {"type": "content", "subtype": "table",
                 "title": "Revenue by Region",
                 "table": {"headers": ["Region", "Q4", "YoY", "vs Target"],
                           "rows": [
                               ["North America", "$2.8M", "+15%", "Met"],
                               ["Europe",        "$1.5M", "+22%", "Exceeded"],
                               ["APAC",          "$0.6M", "+85%", "Exceeded"],
                               ["LATAM",         "$0.3M", "+5%",  "Below"]]}},
                {"type": "section", "number": "03", "title": "Roadmap"},
                {"type": "content", "subtype": "comparison",
                 "title": "Q4 2024 vs Q1 2025",
                 "left":  {"title": "Q4 2024",
                           "points": ["Revenue $5.2M", "Team 58", "NPS 62"]},
                 "right": {"title": "Q1 2025 (target)",
                           "points": ["Revenue $6.0M", "Team 70",
                                      "NPS 65", "v3.0 launched"]}},
                {"type": "summary", "title": "Key Takeaways",
                 "takeaways": ["Strong revenue + team growth",
                               "APAC expansion validated",
                               "Q1 focus: profitability & v3.0",
                               "Hiring plan locked"],
                 "contact": "reports@company.com"}
            ]
        }
    },

    "training": {
        "_desc": "8-page training / education deck (intro + concepts + examples + Q&A).",
        "outline": {
            "title": "Python Async Programming",
            "subtitle": "From callbacks to async/await",
            "meta": "Engineering Training · Module 4",
            "palette": "education",
            "style": "soft",
            "slides": [
                {"type": "toc", "title": "What We'll Cover",
                 "sections": ["Why async?", "asyncio basics",
                              "Common patterns", "Pitfalls", "Q&A"]},
                {"type": "section", "number": "01", "title": "Why async?"},
                {"type": "content", "subtype": "stats",
                 "title": "The IO Wait Problem",
                 "stats": [
                     {"number": "90%", "label": "Time idle", "desc": "in IO-bound app"},
                     {"number": "100x", "label": "Concurrent",
                      "desc": "vs threads"},
                     {"number": "1MB",  "label": "Per coroutine",
                      "desc": "vs 8MB thread"}]},
                {"type": "section", "number": "02", "title": "Basics"},
                {"type": "content", "subtype": "text",
                 "title": "Three Concepts",
                 "bullets": ["Coroutine: function defined with `async def`",
                             "await: pause and yield to event loop",
                             "Event loop: scheduler that runs ready coroutines"]},
                {"type": "content", "subtype": "comparison",
                 "title": "Sync vs Async",
                 "left":  {"title": "Sync",
                           "points": ["1 thread = 1 task at a time",
                                      "Blocks on IO",
                                      "Easy to reason about"]},
                 "right": {"title": "Async",
                           "points": ["1 thread = N concurrent tasks",
                                      "Yields on await",
                                      "Needs careful state mgmt"]}},
                {"type": "section", "number": "03", "title": "Pitfalls"},
                {"type": "content", "subtype": "text",
                 "title": "Top 4 Mistakes",
                 "bullets": ["Calling sync IO inside async func (blocks loop)",
                             "Forgetting to await",
                             "Mixing asyncio + threading carelessly",
                             "Not using semaphores for fan-out"]},
                {"type": "summary", "title": "Takeaways & Q&A",
                 "takeaways": ["Use async for IO-bound, not CPU-bound",
                               "Always await coroutines",
                               "Profile before optimizing",
                               "asyncio.gather + Semaphore = fan-out"]}
            ]
        }
    }
}


# ===========================================================================
# Tail-output convention: every command finishes with [OK]/[ERROR] + the full
# SKILL.md body (frontmatter + first H1 stripped). The model reads SKILL.md as
# the single source of truth for what to do next, so the helper script does
# NOT presume to suggest a specific next step; that judgment belongs to the
# task-level reasoner. _next_step is kept as a no-op for source-compat with
# old call sites scattered across cmd_* functions.
# ===========================================================================
def _next_step(lines):
    return  # intentionally silent; the SKILL.md tail is the only manual.


_SKILL_MD_BODY_CACHE = None


def _load_skill_md_body():
    """Read sibling SKILL.md, strip YAML frontmatter and the first H1 title.
    Cached for the lifetime of the process. Returns "" on any IO error."""
    global _SKILL_MD_BODY_CACHE
    if _SKILL_MD_BODY_CACHE is not None:
        return _SKILL_MD_BODY_CACHE
    body = ""
    try:
        here = os.path.dirname(os.path.abspath(__file__))
        md_path = os.path.normpath(os.path.join(here, os.pardir, "SKILL.md"))
        with open(md_path, "r", encoding="utf-8") as f:
            text = f.read()
        lines = text.splitlines()
        i, n = 0, len(lines)
        # 1) skip YAML frontmatter delimited by leading `---`
        if i < n and lines[i].strip() == "---":
            i += 1
            while i < n and lines[i].strip() != "---":
                i += 1
            if i < n:
                i += 1
        # 2) skip blank lines after frontmatter
        while i < n and not lines[i].strip():
            i += 1
        # 3) skip the first H1 title line (e.g. "# Pptx-Editor")
        if i < n and lines[i].lstrip().startswith("# "):
            i += 1
        # 4) skip blanks after the title
        while i < n and not lines[i].strip():
            i += 1
        body = "\n".join(lines[i:]).rstrip() + "\n"
    except Exception:
        body = ""
    _SKILL_MD_BODY_CACHE = body
    return body


def _print_commands_menu(stream=None):
    """Print SKILL.md (without frontmatter/title) so the model always has the
    full skill manual at the tail of every command's output. Called on both
    success and failure paths."""
    if stream is None:
        stream = sys.stdout
    body = _load_skill_md_body()
    stream.write("\n=== SKILL.md ===\n")
    if body:
        stream.write(body)
    else:
        stream.write("(SKILL.md not found beside this script; "
                     "open skills/pptx-editor/SKILL.md for the manual)\n")


def cmd_examples(args):
    """List examples, print one, or write one to a file."""
    name = args.name
    if not name or name == "list":
        print("Available outline examples (use: examples <name> [--out FILE]):")
        for k, v in EXAMPLES_PPTX.items():
            print(f"  {k:<14}  {v['_desc']}")
        _next_step([
            "Pick a template and dump its outline JSON to a file:",
            "  python pptx_helper.py examples <name> --out ${WORKSPACE}/outline.json",
            "Or take the fast path (no outline editing):",
            "  python pptx_helper.py create --template <name> --out ${WORKSPACE}/deck.pptx",
        ])
        return
    if name not in EXAMPLES_PPTX:
        raise ValueError(f"unknown example '{name}'. choices: "
                         f"{list(EXAMPLES_PPTX)} or 'list'")
    outline_json = json.dumps(EXAMPLES_PPTX[name]["outline"],
                              ensure_ascii=False, indent=2)
    out_path = getattr(args, "out", None)
    if out_path:
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(outline_json)
        print(f"[OK] example '{name}' outline written to {out_path}")
        print(f"[WARN] This file is the UNMODIFIED '{name}' template placeholder. "
              f"If you pass it straight to new-from-outline, the build will be REJECTED.")
        _next_step([
            f"REQUIRED next: overwrite {out_path} with your real task content using write_file.",
            "  - Edit 'title', 'subtitle', and every slides[i].title plus its content fields",
            "    (bullets / rows / stats / items). Keep 'type' and 'subtype' unchanged.",
            f"Then validate: python pptx_helper.py new-from-outline {out_path} --dry-run",
            f"Finally build: python pptx_helper.py new-from-outline {out_path} --out ${{WORKSPACE}}/deck.pptx",
        ])
    else:
        print(outline_json)
        _next_step([
            "Save the JSON above to a file, or re-run with --out FILE to write directly.",
            "Then validate/build with new-from-outline.",
        ])


def _load_json(path):
    """Load JSON tolerating UTF-8 BOM (common when written by Windows tools)."""
    with open(path, "r", encoding="utf-8-sig") as f:
        return json.load(f)


def _outline_signature(outline):
    """Stable content hash of an outline, used to detect unmodified templates.

    Normalized via sort_keys so re-serialization order cannot mask a match.
    """
    return json.dumps(outline, sort_keys=True, ensure_ascii=False)


_EXAMPLE_SIGNATURES = {
    name: _outline_signature(v["outline"]) for name, v in EXAMPLES_PPTX.items()
}


def _detect_unmodified_template(outline):
    """If outline matches any built-in example byte-for-byte, return its name.

    This catches the common LLM mistake: call `examples NAME --out F.json` and
    then pass F.json straight to `new-from-outline` without editing -> the deck
    would contain the template's placeholder copy (e.g. the 'training' example
    is a Python tutorial deck) instead of the user's real task content.
    """
    sig = _outline_signature(outline)
    for name, ref_sig in _EXAMPLE_SIGNATURES.items():
        if sig == ref_sig:
            return name
    return None


def _build_from_outline(outline, out_path, palette_arg=None, style_arg=None):
    """Core deck builder; shared by new-from-outline and create.

    Returns (pal_name, style_name, slide_count).
    """
    pal_name = palette_arg or outline.get("palette", "business")
    style_name = style_arg or outline.get("style", "soft")
    pal = PALETTES.get(pal_name)
    if pal is None:
        if isinstance(pal_name, dict):
            pal = pal_name
        else:
            raise ValueError(f"unknown palette: {pal_name}. choices: {list(PALETTES)}")
    style = STYLES.get(style_name)
    if style is None:
        raise ValueError(f"unknown style: {style_name}. choices: {list(STYLES)}")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]  # blank

    slides_data = list(outline.get("slides", []))
    if outline.get("title") and (not slides_data or slides_data[0].get("type") != "cover"):
        slides_data.insert(0, {
            "type": "cover",
            "title": outline["title"],
            "subtitle": outline.get("subtitle", ""),
            "meta": outline.get("meta", ""),
        })

    page_num = 0
    for sd in slides_data:
        slide = prs.slides.add_slide(blank)
        t = sd.get("type", "content")
        if t == "content":
            sub = sd.get("subtype", "text")
            key = f"content.{sub}"
        else:
            key = t
        builder = _BUILDERS.get(key)
        if builder is None:
            builder = _build_content_text
            sd = {"title": sd.get("title", f"[unknown type: {t}]"),
                  "bullets": sd.get("bullets", [])}
        if key == "cover":
            builder(slide, sd, pal, style)
        else:
            page_num += 1
            builder(slide, sd, pal, style, page_num)

    prs.save(out_path)
    return pal_name, style_name, len(slides_data)


def cmd_new_from_outline(args):
    outline_path = args.outline_pos or args.outline_opt
    if not outline_path:
        raise ValueError("outline path is required "
                         "(pass as positional or --outline PATH)")
    outline = _load_json(outline_path)

    # Guard: refuse to build when the outline is byte-for-byte identical to a
    # built-in example (the LLM forgot to edit it). Without this check the
    # generated deck would contain the template's placeholder copy (e.g. the
    # 'training' example is a Python tutorial deck), not the user's task
    # content. Signal as validation failure so the error surfaces in --dry-run
    # and in the real build path.
    tmpl = _detect_unmodified_template(outline)
    if tmpl:
        print(f"[ERROR] outline.json is the UNMODIFIED '{tmpl}' example template.")
        print(f"        Building now would produce a generic '{tmpl}' deck, NOT your task content.")
        print(f"        You MUST edit the following fields to match the task before re-running:")
        print(f"          - top-level: 'title', 'subtitle'")
        print(f"          - each slides[i]: 'title' and content fields "
              f"(bullets / rows / stats / items depending on 'type'/'subtype')")
        print(f"        Use write_file to overwrite {outline_path} with your customized outline,")
        print(f"        then re-run `new-from-outline`.")
        _next_step([
            f"Overwrite {outline_path} with real task content (keep 'type'/'subtype' unchanged),",
            f"then validate: python pptx_helper.py new-from-outline {outline_path} --dry-run",
        ])
        sys.exit(1)

    # Always run schema/density validation; print issues even when not dry-run.
    issues = _validate_outline(outline)
    if issues:
        print("[VALIDATION] issues found:")
        for it in issues:
            print(f"  - {it}")
    else:
        print("[VALIDATION] OK")

    if args.dry_run:
        print("[DRY-RUN] no file written.")
        if issues:
            _next_step([
                "Fix the issues listed above in the outline JSON, then re-run:",
                f"  python pptx_helper.py new-from-outline {outline_path} --dry-run",
            ])
            sys.exit(1)
        _next_step([
            "Outline passed validation. To actually build the file, run:",
            f"  python pptx_helper.py new-from-outline {outline_path} --out ${{WORKSPACE}}/deck.pptx",
        ])
        return

    if not args.out:
        raise ValueError("--out PATH is required unless --dry-run is set")
    pal_name, style_name, n = _build_from_outline(
        outline, args.out, args.palette, args.style)
    print(f"[OK] new-from-outline -> {args.out} "
          f"(palette={pal_name}, style={style_name}, slides={n})")
    _next_step([
        f"Deck generated at {args.out}",
        f"Inspect layout:  python pptx_helper.py inspect {args.out}",
        f"Tweak text:      python pptx_helper.py set-text {args.out} "
        "--slide N --shape K --text \"...\"",
    ])


def cmd_create(args):
    """Fast path: generate a pptx directly from a built-in template."""
    if args.outline_file:
        outline = _load_json(args.outline_file)
        src = f"outline file {args.outline_file}"
    else:
        if args.template not in EXAMPLES_PPTX:
            raise ValueError(f"unknown template '{args.template}'. "
                             f"choices: {list(EXAMPLES_PPTX)}")
        outline = deepcopy(EXAMPLES_PPTX[args.template]["outline"])
        src = f"built-in template '{args.template}'"
    issues = _validate_outline(outline)
    if issues:
        print("[VALIDATION] issues:")
        for it in issues:
            print(f"  - {it}")
    else:
        print("[VALIDATION] OK")
    pal_name, style_name, n = _build_from_outline(
        outline, args.out, args.palette, args.style)
    print(f"[OK] create -> {args.out} (from {src}; "
          f"palette={pal_name}, style={style_name}, slides={n})")
    _next_step([
        f"Deck generated at {args.out}",
        f"Inspect:   python pptx_helper.py inspect {args.out}",
        f"Edit text: python pptx_helper.py set-text {args.out} "
        "--slide N --shape K --text \"...\"",
        "Custom content path: examples <name> --out outline.json "
        "-> edit fields -> new-from-outline <file> --out <deck.pptx>",
    ])


# ===========================================================================
# Plan-first API (RECOMMENDED): closed-set layout catalog with strict slot
# schema. The model picks (theme, layout) from a small enumerated list and
# fills named slots; everything else (colors, fonts, sizing, badges) is
# locked down so the model has zero design freedom and the deck stays
# visually consistent.
#
# Workflow:
#   1) python pptx_helper.py catalog            (one-shot view of choices)
#   2) write_file plan.json with theme + slides
#   3) python pptx_helper.py plan plan.json --dry-run   (validate slots)
#   4) python pptx_helper.py plan plan.json --out F.pptx
# ===========================================================================

# Slot-type vocabulary used in LAYOUTS:
#   "str"             - non-empty string
#   "str?"            - optional string
#   "list[str]"       - list of strings (with min/max items)
#   "list[list[str]]" - 2D list (table rows)
#   "obj"             - single dict with named sub-slots (item_slots)
#   "list[obj]"       - list of dicts with named sub-slots
#   "path"            - filesystem path string (existence checked, soft-warn)

LAYOUTS = {
    # --- Front matter ---
    "cover": {
        "role": "cover",
        "summary": "Title slide with big centered title.",
        "slots": {
            "title":    {"type": "str",  "required": True,  "max_chars": 60},
            "subtitle": {"type": "str?", "required": False, "max_chars": 80},
            "meta":     {"type": "str?", "required": False, "max_chars": 60},
        },
        "builder": ("cover", {}),
    },
    "toc": {
        "role": "toc",
        "summary": "Table of contents with 3-6 numbered sections.",
        "slots": {
            "title": {"type": "str?", "required": False, "max_chars": 40},
            "items": {"type": "list[obj]", "required": True, "min": 3, "max": 6,
                      "item_slots": {
                          "title": {"type": "str", "required": True, "max_chars": 40},
                          "desc":  {"type": "str?", "required": False, "max_chars": 60},
                      }},
        },
        "builder": ("toc", {}),
    },
    "section": {
        "role": "section",
        "summary": "Section divider with big number + title.",
        "slots": {
            "number": {"type": "str", "required": True, "max_chars": 4},
            "title":  {"type": "str", "required": True, "max_chars": 40},
            "intro":  {"type": "str?", "required": False, "max_chars": 100},
        },
        "builder": ("section", {}),
    },
    # --- Content slides (pure variants) ---
    "text": {
        "role": "content",
        "summary": "Full-width title + 2-7 bulleted lines.",
        "slots": {
            "title":   {"type": "str", "required": True, "max_chars": 50},
            "bullets": {"type": "list[str]", "required": True, "min": 2, "max": 7},
        },
        "builder": ("content.text", {}),
    },
    "table": {
        "role": "content",
        "summary": "Full-width data table (2-5 cols, 1-8 rows).",
        "slots": {
            "title":   {"type": "str", "required": True, "max_chars": 50},
            "headers": {"type": "list[str]", "required": True, "min": 2, "max": 5},
            "rows":    {"type": "list[list[str]]", "required": True,
                        "min_rows": 1, "max_rows": 8},
        },
        "builder": ("content.table", {}),
    },
    "stats": {
        "role": "content",
        "summary": "2-4 large stat callouts (big number + label + optional desc).",
        "slots": {
            "title": {"type": "str", "required": True, "max_chars": 50},
            "stats": {"type": "list[obj]", "required": True, "min": 2, "max": 4,
                      "item_slots": {
                          "value": {"type": "str", "required": True, "max_chars": 8},
                          "label": {"type": "str", "required": True, "max_chars": 20},
                          "desc":  {"type": "str?", "required": False, "max_chars": 30},
                      }},
        },
        "builder": ("content.stats", {}),
    },
    "compare": {
        "role": "content",
        "summary": "Two side-by-side cards (A vs B), each with label + 2-6 bullets.",
        "slots": {
            "title": {"type": "str", "required": True, "max_chars": 50},
            "left":  {"type": "obj", "required": True, "item_slots": {
                "label":  {"type": "str", "required": True, "max_chars": 30},
                "points": {"type": "list[str]", "required": True, "min": 2, "max": 6},
            }},
            "right": {"type": "obj", "required": True, "item_slots": {
                "label":  {"type": "str", "required": True, "max_chars": 30},
                "points": {"type": "list[str]", "required": True, "min": 2, "max": 6},
            }},
        },
        "builder": ("content.comparison", {}),
    },
    "timeline": {
        "role": "content",
        "summary": "3-5 horizontal numbered steps (process / journey).",
        "slots": {
            "title": {"type": "str", "required": True, "max_chars": 50},
            "steps": {"type": "list[obj]", "required": True, "min": 3, "max": 5,
                      "item_slots": {
                          "label": {"type": "str", "required": True, "max_chars": 25},
                          "desc":  {"type": "str?", "required": False, "max_chars": 60},
                      }},
        },
        "builder": ("__timeline__", {}),
    },
    # --- Mixed: text + table positioned variants (the 4 the user asked for) ---
    "text_table_right": {
        "role": "content",
        "summary": "Left bullets (1-5) + right table (2-5 cols, 1-6 rows).",
        "slots": {
            "title":   {"type": "str", "required": True, "max_chars": 50},
            "bullets": {"type": "list[str]", "required": True, "min": 1, "max": 5},
            "headers": {"type": "list[str]", "required": True, "min": 2, "max": 5},
            "rows":    {"type": "list[list[str]]", "required": True,
                        "min_rows": 1, "max_rows": 6},
        },
        "builder": ("__text_table_split__", {"table_pos": "right"}),
    },
    "text_table_left": {
        "role": "content",
        "summary": "Left table + right bullets (1-5).",
        "slots": {
            "title":   {"type": "str", "required": True, "max_chars": 50},
            "bullets": {"type": "list[str]", "required": True, "min": 1, "max": 5},
            "headers": {"type": "list[str]", "required": True, "min": 2, "max": 5},
            "rows":    {"type": "list[list[str]]", "required": True,
                        "min_rows": 1, "max_rows": 6},
        },
        "builder": ("__text_table_split__", {"table_pos": "left"}),
    },
    "text_table_bottom": {
        "role": "content",
        "summary": "Top bullets (1-3) + bottom table (2-5 cols, 1-5 rows).",
        "slots": {
            "title":   {"type": "str", "required": True, "max_chars": 50},
            "bullets": {"type": "list[str]", "required": True, "min": 1, "max": 3},
            "headers": {"type": "list[str]", "required": True, "min": 2, "max": 5},
            "rows":    {"type": "list[list[str]]", "required": True,
                        "min_rows": 1, "max_rows": 5},
        },
        "builder": ("__text_table_split__", {"table_pos": "bottom"}),
    },
    "text_table_top": {
        "role": "content",
        "summary": "Top table + bottom bullets (1-3).",
        "slots": {
            "title":   {"type": "str", "required": True, "max_chars": 50},
            "bullets": {"type": "list[str]", "required": True, "min": 1, "max": 3},
            "headers": {"type": "list[str]", "required": True, "min": 2, "max": 5},
            "rows":    {"type": "list[list[str]]", "required": True,
                        "min_rows": 1, "max_rows": 5},
        },
        "builder": ("__text_table_split__", {"table_pos": "top"}),
    },
    # --- Mixed: text + image positioned variants ---
    "text_image_right": {
        "role": "content",
        "summary": "Left bullets (1-5) + right image (placeholder if path missing).",
        "slots": {
            "title":      {"type": "str", "required": True, "max_chars": 50},
            "bullets":    {"type": "list[str]", "required": True, "min": 1, "max": 5},
            "image_path": {"type": "path", "required": False},
            "caption":    {"type": "str?", "required": False, "max_chars": 40},
        },
        "builder": ("__text_image_split__", {"image_pos": "right"}),
    },
    "text_image_left": {
        "role": "content",
        "summary": "Left image (placeholder if path missing) + right bullets (1-5).",
        "slots": {
            "title":      {"type": "str", "required": True, "max_chars": 50},
            "bullets":    {"type": "list[str]", "required": True, "min": 1, "max": 5},
            "image_path": {"type": "path", "required": False},
            "caption":    {"type": "str?", "required": False, "max_chars": 40},
        },
        "builder": ("__text_image_split__", {"image_pos": "left"}),
    },
    # --- Closing ---
    "summary": {
        "role": "summary",
        "summary": "Key takeaways list (2-5) + optional contact line.",
        "slots": {
            "title":     {"type": "str", "required": True, "max_chars": 40},
            "takeaways": {"type": "list[str]", "required": True, "min": 2, "max": 5},
            "contact":   {"type": "str?", "required": False, "max_chars": 60},
        },
        "builder": ("summary", {}),
    },
}

THEMES = list(PALETTES.keys())  # business / tech / wellness / elegant / education / platinum


def _slot_kind(spec):
    return spec.get("type", "str")


def _validate_slot(name, spec, value, prefix=""):
    """Recursively validate one slot. Return list of error strings."""
    errs = []
    kind = _slot_kind(spec)
    optional = kind.endswith("?")
    base = kind.rstrip("?")
    required = spec.get("required", False) and not optional
    is_empty = (value is None) or (isinstance(value, (str, list, dict)) and len(value) == 0)
    if is_empty:
        if required:
            errs.append(f"{prefix}{name}: missing required slot (type={kind})")
        return errs
    if base == "str":
        if not isinstance(value, str):
            errs.append(f"{prefix}{name}: expected string, got {type(value).__name__}")
        elif "max_chars" in spec and len(value) > spec["max_chars"]:
            errs.append(f"{prefix}{name}: {len(value)} chars > max_chars={spec['max_chars']}")
    elif base == "path":
        if not isinstance(value, str):
            errs.append(f"{prefix}{name}: expected path string")
        elif not os.path.exists(value):
            errs.append(f"{prefix}{name}: file not found '{value}' (will render placeholder)")
    elif base == "list[str]":
        if not isinstance(value, list):
            errs.append(f"{prefix}{name}: expected list of strings"); return errs
        if "min" in spec and len(value) < spec["min"]:
            errs.append(f"{prefix}{name}: {len(value)} items < min={spec['min']}")
        if "max" in spec and len(value) > spec["max"]:
            errs.append(f"{prefix}{name}: {len(value)} items > max={spec['max']}")
        for i, v in enumerate(value):
            if not isinstance(v, str):
                errs.append(f"{prefix}{name}[{i}]: expected string")
    elif base == "list[list[str]]":
        if not isinstance(value, list):
            errs.append(f"{prefix}{name}: expected 2D list"); return errs
        if "min_rows" in spec and len(value) < spec["min_rows"]:
            errs.append(f"{prefix}{name}: {len(value)} rows < min_rows={spec['min_rows']}")
        if "max_rows" in spec and len(value) > spec["max_rows"]:
            errs.append(f"{prefix}{name}: {len(value)} rows > max_rows={spec['max_rows']}")
        for ri, row in enumerate(value):
            if not isinstance(row, list):
                errs.append(f"{prefix}{name}[{ri}]: expected row list")
    elif base == "obj":
        if not isinstance(value, dict):
            errs.append(f"{prefix}{name}: expected object"); return errs
        for sn, ss in spec.get("item_slots", {}).items():
            errs.extend(_validate_slot(sn, ss, value.get(sn), prefix=f"{prefix}{name}."))
    elif base == "list[obj]":
        if not isinstance(value, list):
            errs.append(f"{prefix}{name}: expected list of objects"); return errs
        if "min" in spec and len(value) < spec["min"]:
            errs.append(f"{prefix}{name}: {len(value)} items < min={spec['min']}")
        if "max" in spec and len(value) > spec["max"]:
            errs.append(f"{prefix}{name}: {len(value)} items > max={spec['max']}")
        for i, item in enumerate(value):
            if not isinstance(item, dict):
                errs.append(f"{prefix}{name}[{i}]: expected object"); continue
            for sn, ss in spec.get("item_slots", {}).items():
                errs.extend(_validate_slot(sn, ss, item.get(sn),
                                           prefix=f"{prefix}{name}[{i}]."))
    return errs


def _validate_plan(plan):
    """Return list of issue strings; empty list = OK."""
    issues = []
    if not isinstance(plan, dict):
        return ["plan must be a JSON object"]
    theme = plan.get("theme", "business")
    if theme not in PALETTES:
        issues.append(f"theme '{theme}' unknown; choices: {THEMES}")
    slides = plan.get("slides", [])
    if not isinstance(slides, list) or not slides:
        return issues + ["'slides' must be a non-empty array"]
    for i, sd in enumerate(slides):
        prefix = f"slides[{i}]: "
        if not isinstance(sd, dict):
            issues.append(f"{prefix}must be an object"); continue
        layout = sd.get("layout")
        if not layout or layout not in LAYOUTS:
            issues.append(f"{prefix}layout='{layout}' not in catalog. "
                          f"Run `pptx_helper.py catalog` for valid IDs.")
            continue
        spec = LAYOUTS[layout]
        slots = sd.get("slots", {})
        if not isinstance(slots, dict):
            issues.append(f"{prefix}'slots' must be an object"); continue
        for sn, ss in spec["slots"].items():
            issues.extend(_validate_slot(sn, ss, slots.get(sn), prefix=prefix))
        # Reject unknown slot keys (catches typos / hallucinated fields)
        unknown = set(slots) - set(spec["slots"])
        if unknown:
            issues.append(f"{prefix}unknown slot(s): {sorted(unknown)} "
                          f"(valid: {list(spec['slots'])})")
    return issues


def _slots_to_legacy_data(layout_id, slots):
    """Translate slot dict to the dict shape expected by legacy _build_* funcs."""
    if layout_id == "toc":
        return {
            "title": slots.get("title", "Table of Contents"),
            "sections": slots.get("items", []),
        }
    if layout_id == "stats":
        return {
            "title": slots.get("title", ""),
            "stats": [{"number": s.get("value", ""),
                       "label": s.get("label", ""),
                       "desc": s.get("desc", "")}
                      for s in slots.get("stats", [])],
        }
    if layout_id == "compare":
        out = {"title": slots.get("title", "")}
        for k in ("left", "right"):
            sd = slots.get(k, {}) or {}
            out[k] = {"title": sd.get("label", ""),
                      "points": sd.get("points", [])}
        return out
    if layout_id == "table":
        return {
            "title": slots.get("title", ""),
            "table": {"headers": slots.get("headers", []),
                      "rows": slots.get("rows", [])},
        }
    if layout_id == "summary":
        return {
            "title": slots.get("title", ""),
            "takeaways": slots.get("takeaways", []),
            "contact": slots.get("contact", ""),
        }
    # cover / section / text / mixed-split / timeline pass through directly
    return dict(slots)


def _build_from_plan(plan, out_path):
    """Build deck from a plan dict. Returns (theme_name, slide_count)."""
    theme = plan.get("theme", "business")
    pal = PALETTES.get(theme, PALETTES["business"])
    style = STYLES["soft"]  # Style is fixed; theme drives all visuals.
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    page_num = 0
    total = 0
    for sd in plan.get("slides", []):
        layout_id = sd["layout"]
        spec = LAYOUTS[layout_id]
        legacy = _slots_to_legacy_data(layout_id, sd.get("slots", {}))
        slide = prs.slides.add_slide(blank)
        builder_key, extra = spec["builder"]
        if builder_key == "cover":
            _build_cover(slide, legacy, pal, style)
        elif builder_key == "__timeline__":
            page_num += 1
            _build_content_timeline(slide, legacy, pal, style, page_num)
        elif builder_key == "__text_table_split__":
            page_num += 1
            _build_text_table_split(slide, legacy, pal, style, page_num,
                                    table_pos=extra["table_pos"])
        elif builder_key == "__text_image_split__":
            page_num += 1
            _build_text_image_split(slide, legacy, pal, style, page_num,
                                    image_pos=extra["image_pos"])
        else:
            page_num += 1
            _BUILDERS[builder_key](slide, legacy, pal, style, page_num)
        total += 1
    prs.save(out_path)
    return theme, total


# ===========================================================================
# Direct-CLI deck builder: one subcommand per slide. No JSON file, no schema,
# no validation step. Each add-* opens the deck, appends one slide, writes it
# back. Out-of-bounds inputs are auto-clipped + emit [WARN] (never fail).
#
# Flow (the model only memorizes this):
#   pptx_helper.py new-deck FILE.pptx --theme business
#   pptx_helper.py add-cover FILE.pptx --title "..." --subtitle "..."
#   pptx_helper.py add-table FILE.pptx --title "..." --headers "h1,h2,h3"
#                                       --row "v1,v2,v3" --row "v4,v5,v6"
#   ...
# ===========================================================================
def _resolve_theme(prs):
    """Return (theme_name, palette, style) using deck metadata."""
    meta = _meta_load(prs)
    theme = meta["theme"] if meta["theme"] in PALETTES else "business"
    return theme, PALETTES[theme], STYLES["soft"], meta


def _commit_slide(prs, args, kind, meta, *, increment_page=True):
    """Save deck after adding a slide and print [OK] / NEXT_STEP.

    `meta` is the CALLER'S mutated copy. We must NOT reload from prs here
    (caller may have updated section_count / theme), otherwise local edits
    get clobbered. _meta_save writes it back unconditionally.
    """
    if increment_page:
        meta["next_page_num"] += 1
    _meta_save(prs, meta)
    prs.save(args.file)
    n = len(prs.slides)
    print(f"[OK] {kind} -> {args.file} (slides={n}, theme={meta['theme']})")
    _next_step([
        f"Slide #{n} added. Continue with another add-* (e.g. add-text, add-table, add-summary)",
        f"or finish: python pptx_helper.py inspect {args.file}",
    ])


def cmd_new_deck(args):
    theme = (args.theme or "business").strip().lower()
    if theme not in PALETTES:
        print(f"[WARN] theme '{theme}' unknown; falling back to 'business' "
              f"(valid: {list(PALETTES)})")
        theme = "business"

    # --template PATH: load an existing .pptx/.potx as the base deck so we
    # inherit its slide master (logos, brand fonts, master colors, footer).
    # Any slides the template already contains are KEPT and our add-* commands
    # append AFTER them. This lets users plug in a corp template's intro/cover
    # and just have us fill in the data pages.
    tpl = getattr(args, "template", None)
    used_template = False
    if tpl:
        if not os.path.isfile(tpl):
            raise FileNotFoundError(
                f"new-deck: --template '{tpl}' not found. "
                f"Pass an existing .pptx or .potx path, or omit --template "
                f"to start from a blank deck."
            )
        prs = Presentation(tpl)
        used_template = True
    else:
        prs = Presentation()

    # Our add-* layouts use absolute Inches positioning against a 16:9 canvas
    # (SLIDE_W_IN x SLIDE_H_IN). If a custom template ships with a different
    # canvas (e.g. 4:3 10x7.5), warn loudly because content boxes will leave
    # gaps or overflow on those template's existing slides. We still force the
    # size override so OUR newly-appended slides render correctly.
    if used_template:
        tpl_w = float(prs.slide_width) / 914400.0   # EMU -> inches
        tpl_h = float(prs.slide_height) / 914400.0
        if abs(tpl_w - SLIDE_W_IN) > 0.05 or abs(tpl_h - SLIDE_H_IN) > 0.05:
            print(f"[WARN] template canvas {tpl_w:.2f}x{tpl_h:.2f}in differs "
                  f"from our 16:9 layout {SLIDE_W_IN:.2f}x{SLIDE_H_IN:.2f}in; "
                  f"forcing 16:9. Template's existing slides may look off.")
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # Continue page-number footer numbering AFTER any slides the template
    # already had. _add_page_badge reads this counter when rendering badges.
    n_existing = len(prs.slides)
    _meta_save(prs, {"theme": theme,
                     "next_page_num": max(1, n_existing + 1),
                     "section_count": 0})
    prs.save(args.file)
    if used_template:
        print(f"[OK] new-deck -> {args.file} (theme={theme}, "
              f"template={tpl}, kept {n_existing} template slide(s), "
              f"next_page={n_existing + 1})")
    else:
        print(f"[OK] new-deck -> {args.file} (theme={theme}, slides=0)")


def cmd_add_cover(args):
    prs = _open_deck(args.file)
    _, pal, style, _meta = _resolve_theme(prs)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "title":    _clip_str(args.title, _LIM["title"], where="title"),
        "subtitle": _clip_str(args.subtitle or "", _LIM["subtitle"], where="subtitle"),
        "meta":     _clip_str(args.meta or "", _LIM["meta"], where="meta"),
    }
    _build_cover(slide, data, pal, style)
    # Cover does NOT consume a page number (no badge).
    _commit_slide(prs, args, "add-cover", _meta, increment_page=False)


def cmd_add_toc(args):
    prs = _open_deck(args.file)
    _, pal, style, meta = _resolve_theme(prs)
    items_raw = list(args.item or [])
    items_raw = _clip_list(items_raw, _LIM["toc_items_max"], where="toc items")
    if len(items_raw) < _LIM["toc_items_min"]:
        print(f"[WARN] toc has only {len(items_raw)} items "
              f"(recommended >= {_LIM['toc_items_min']}); rendering anyway")
    sections = []
    for it in items_raw:
        title, desc = _split_kv(it, sep=";", n_fields=2)
        sections.append({
            "title": _clip_str(title, _LIM["toc_item_title"], where="toc item title"),
            "desc":  _clip_str(desc,  _LIM["toc_item_desc"],  where="toc item desc"),
        })
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "title": _clip_str(args.title or "Table of Contents", _LIM["title"], where="toc title"),
        "sections": sections,
    }
    _build_toc(slide, data, pal, style, meta["next_page_num"])
    _commit_slide(prs, args, "add-toc", meta)


def cmd_add_section(args):
    prs = _open_deck(args.file)
    _, pal, style, meta = _resolve_theme(prs)
    meta["section_count"] += 1
    number = args.number if args.number else f"{meta['section_count']:02d}"
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "number": _clip_str(number, _LIM["section_number"], where="section number"),
        "title":  _clip_str(args.title, _LIM["title"], where="section title"),
        "intro":  _clip_str(args.intro or "", _LIM["section_intro"], where="section intro"),
    }
    _build_section(slide, data, pal, style, meta["next_page_num"])
    _commit_slide(prs, args, "add-section", meta)


def _normalize_bullets(raw, *, max_n, min_n, where):
    bullets = [b for b in (raw or []) if str(b).strip()]
    bullets = _clip_list(bullets, max_n, where=where)
    bullets = [_clip_str(b, _LIM["bullet"], where=f"{where}[{i}]") for i, b in enumerate(bullets)]
    if len(bullets) < min_n:
        print(f"[WARN] {where}: only {len(bullets)} provided (recommended >= {min_n}); rendering anyway")
    return bullets


def cmd_add_text(args):
    prs = _open_deck(args.file)
    _, pal, style, meta = _resolve_theme(prs)
    bullets = _normalize_bullets(args.bullet,
                                 max_n=_LIM["bullets_max"],
                                 min_n=_LIM["bullets_min"],
                                 where="bullets")
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "title":   _clip_str(args.title, _LIM["title"], where="text title"),
        "bullets": bullets,
    }
    _build_content_text(slide, data, pal, style, meta["next_page_num"])
    _commit_slide(prs, args, "add-text", meta)


def _normalize_table(headers_raw, rows_raw):
    """Parse + clip headers and rows; pad/trim each row to header count."""
    headers = _split_csv(headers_raw, sep=",")
    headers = _clip_list(headers, _LIM["headers_max"], where="table headers")
    headers = [_clip_str(h, _LIM["header_cell"], where=f"header[{i}]") for i, h in enumerate(headers)]
    nc = len(headers)
    rows = []
    for ri, r in enumerate(rows_raw or []):
        cells = _split_csv(r, sep=",")
        # Pad/trim row to nc columns
        if len(cells) < nc:
            cells = cells + [""] * (nc - len(cells))
        elif len(cells) > nc:
            print(f"[WARN] row[{ri}]: kept first {nc} of {len(cells)} cells")
            cells = cells[:nc]
        cells = [_clip_str(c, _LIM["table_cell"], where=f"row[{ri}][{ci}]") for ci, c in enumerate(cells)]
        rows.append(cells)
    rows = _clip_list(rows, _LIM["rows_max"], where="table rows")
    return headers, rows


def cmd_add_table(args):
    """Append a headed table slide.

    LLM-friendly invocation forms accepted (positional aligns with
    xlsx add-sheet so the model does not have to remember two styles):
        add-table FILE.pptx --title T --headers "a,b" --row "1,2" --row "3,4"  (canonical)
        add-table FILE.pptx "T" "a,b" "1,2" "3,4"                              (positional fallback)
    """
    # Merge keyword args with positional fallback in slot order: TITLE, HEADERS, ROW...
    pos = list(getattr(args, "pos_args", []) or [])
    title = args.title
    headers_str = args.headers
    if title is None and pos:
        title = pos.pop(0)
    if headers_str is None and pos:
        headers_str = pos.pop(0)
    rows_raw = list(args.row or []) + [s for s in pos if str(s).strip()]
    if not title or not str(title).strip():
        sys.stderr.write("[ERROR] add-table: missing title. Use one of:\n"
                         "  add-table FILE.pptx --title \"T\" --headers \"a,b\" --row \"1,2\"\n"
                         "  add-table FILE.pptx \"T\" \"a,b\" \"1,2\" \"3,4\"\n")
        _print_commands_menu(stream=sys.stderr)
        sys.exit(2)
    if not headers_str or not str(headers_str).strip():
        sys.stderr.write("[ERROR] add-table: missing headers. Use one of:\n"
                         "  add-table FILE.pptx --title \"T\" --headers \"a,b\" --row \"1,2\"\n"
                         "  add-table FILE.pptx \"T\" \"a,b\" \"1,2\" \"3,4\"\n")
        _print_commands_menu(stream=sys.stderr)
        sys.exit(2)
    prs = _open_deck(args.file)
    _, pal, style, meta = _resolve_theme(prs)
    headers, rows = _normalize_table(headers_str, rows_raw)
    if not headers:
        print("[WARN] table has no headers; rendering empty placeholder")
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "title": _clip_str(title, _LIM["title"], where="table title"),
        "table": {"headers": headers, "rows": rows},
    }
    _build_content_table(slide, data, pal, style, meta["next_page_num"])
    _commit_slide(prs, args, "add-table", meta)


def cmd_add_stats(args):
    prs = _open_deck(args.file)
    _, pal, style, meta = _resolve_theme(prs)
    raw = list(args.stat or [])
    raw = _clip_list(raw, _LIM["stats_max"], where="stats")
    stats = []
    for i, s in enumerate(raw):
        value, label, desc = _split_kv(s, sep=";", n_fields=3)
        stats.append({
            "number": _clip_str(value, _LIM["stat_value"], where=f"stat[{i}].value"),
            "label":  _clip_str(label, _LIM["stat_label"], where=f"stat[{i}].label"),
            "desc":   _clip_str(desc,  _LIM["stat_desc"],  where=f"stat[{i}].desc"),
        })
    if len(stats) < _LIM["stats_min"]:
        print(f"[WARN] stats: only {len(stats)} provided "
              f"(recommended >= {_LIM['stats_min']}); rendering anyway")
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "title": _clip_str(args.title, _LIM["title"], where="stats title"),
        "stats": stats,
    }
    _build_content_stats(slide, data, pal, style, meta["next_page_num"])
    _commit_slide(prs, args, "add-stats", meta)


def cmd_add_compare(args):
    prs = _open_deck(args.file)
    _, pal, style, meta = _resolve_theme(prs)
    left_pts = _normalize_bullets(args.left_point,
                                  max_n=_LIM["compare_points_max"],
                                  min_n=_LIM["compare_points_min"],
                                  where="left points")
    right_pts = _normalize_bullets(args.right_point,
                                   max_n=_LIM["compare_points_max"],
                                   min_n=_LIM["compare_points_min"],
                                   where="right points")
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "title": _clip_str(args.title, _LIM["title"], where="compare title"),
        "left":  {
            "title":  _clip_str(args.left_title or "A", _LIM["compare_col_title"], where="left.title"),
            "points": [_clip_str(p, _LIM["compare_point"], where=f"left.points[{i}]")
                        for i, p in enumerate(left_pts)],
        },
        "right": {
            "title":  _clip_str(args.right_title or "B", _LIM["compare_col_title"], where="right.title"),
            "points": [_clip_str(p, _LIM["compare_point"], where=f"right.points[{i}]")
                        for i, p in enumerate(right_pts)],
        },
    }
    _build_content_comparison(slide, data, pal, style, meta["next_page_num"])
    _commit_slide(prs, args, "add-compare", meta)


def cmd_add_timeline(args):
    prs = _open_deck(args.file)
    _, pal, style, meta = _resolve_theme(prs)
    raw = list(args.step or [])
    raw = _clip_list(raw, _LIM["steps_max"], where="timeline steps")
    steps = []
    for i, s in enumerate(raw):
        label, desc = _split_kv(s, sep=";", n_fields=2)
        steps.append({
            "label": _clip_str(label, _LIM["step_label"], where=f"step[{i}].label"),
            "desc":  _clip_str(desc,  _LIM["step_desc"],  where=f"step[{i}].desc"),
        })
    if len(steps) < _LIM["steps_min"]:
        print(f"[WARN] timeline: only {len(steps)} steps "
              f"(recommended >= {_LIM['steps_min']}); rendering anyway")
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "title": _clip_str(args.title, _LIM["title"], where="timeline title"),
        "steps": steps,
    }
    _build_content_timeline(slide, data, pal, style, meta["next_page_num"])
    _commit_slide(prs, args, "add-timeline", meta)


def cmd_add_text_table(args):
    prs = _open_deck(args.file)
    _, pal, style, meta = _resolve_theme(prs)
    pos = (args.pos or "right").lower()
    if pos not in ("left", "right", "top", "bottom"):
        print(f"[WARN] --pos '{pos}' not in [left,right,top,bottom]; using 'right'")
        pos = "right"
    bullets = _normalize_bullets(args.bullet, max_n=5, min_n=1, where="bullets")
    headers, rows = _normalize_table(args.headers, args.row)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "title":   _clip_str(args.title, _LIM["title"], where="text-table title"),
        "bullets": bullets,
        "headers": headers,
        "rows":    rows,
    }
    _build_text_table_split(slide, data, pal, style, meta["next_page_num"], table_pos=pos)
    _commit_slide(prs, args, "add-text-table", meta)


def cmd_add_text_image(args):
    prs = _open_deck(args.file)
    _, pal, style, meta = _resolve_theme(prs)
    pos = (args.pos or "right").lower()
    if pos not in ("left", "right"):
        print(f"[WARN] --pos '{pos}' not in [left,right]; using 'right'")
        pos = "right"
    bullets = _normalize_bullets(args.bullet, max_n=5, min_n=1, where="bullets")
    if args.image and not os.path.exists(args.image):
        print(f"[WARN] image not found: {args.image}; will render placeholder")
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "title":      _clip_str(args.title, _LIM["title"], where="text-image title"),
        "bullets":    bullets,
        "image_path": args.image or "",
        "caption":    _clip_str(args.caption or "", _LIM["caption"], where="caption"),
    }
    _build_text_image_split(slide, data, pal, style, meta["next_page_num"], image_pos=pos)
    _commit_slide(prs, args, "add-text-image", meta)


def cmd_add_summary(args):
    prs = _open_deck(args.file)
    _, pal, style, meta = _resolve_theme(prs)
    raw = list(args.takeaway or [])
    raw = _clip_list(raw, _LIM["takeaways_max"], where="takeaways")
    takeaways = [_clip_str(t, _LIM["takeaway"], where=f"takeaway[{i}]")
                  for i, t in enumerate(raw) if str(t).strip()]
    if len(takeaways) < _LIM["takeaways_min"]:
        print(f"[WARN] summary: only {len(takeaways)} takeaways "
              f"(recommended >= {_LIM['takeaways_min']}); rendering anyway")
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    data = {
        "title":     _clip_str(args.title or "Key Takeaways", _LIM["title"], where="summary title"),
        "takeaways": takeaways,
        "contact":   _clip_str(args.contact or "", _LIM["contact"], where="contact"),
    }
    _build_summary(slide, data, pal, style, meta["next_page_num"])
    _commit_slide(prs, args, "add-summary", meta)


def cmd_catalog(args):
    """Cookbook: themes + add-* commands + a working full-deck example.

    The model is expected to read this ONCE at task start, then run
    `new-deck` followed by a sequence of `add-*` commands. No JSON file,
    no schema, no dry-run — every parameter is auto-clipped if it
    overflows the slide budget (loud [WARN] but never fail).
    """
    print("=== STEP 1: pick a theme (controls all colors/fonts) ===")
    print("  business    deep navy + red       (corporate / quarterly)")
    print("  tech        deep blue + amber     (product launch / AI)")
    print("  wellness    teal + peach          (health / education / kids)")
    print("  elegant     plum + grey           (luxury / consulting / culture)")
    print("  education   teal + orange         (training / academic)")
    print("  platinum    black + gold          (finance / premium brand)")
    print()
    print("=== STEP 2: init an empty deck (creates FILE.pptx) ===")
    print("  python pptx_helper.py new-deck FILE.pptx --theme business")
    print()
    print("=== STEP 3: append slides (one command per page) ===")
    print()
    print("  --- add-cover -------------------------------------------------")
    print("  Big title page. Use ONCE at the very beginning.")
    print("    pptx_helper.py add-cover FILE.pptx \\")
    print("      --title \"Q1 2026 Flash Price Report\" \\")
    print("      --subtitle \"Weekly snapshot, 2026-04-27\" \\")
    print("      --meta    \"prepared by Research Team\"")
    print()
    print("  --- add-toc ---------------------------------------------------")
    print("  Table-of-contents page (3..6 items; --item repeated)")
    print("    pptx_helper.py add-toc FILE.pptx \\")
    print("      --title \"Agenda\" \\")
    print("      --item  \"Market overview;1-page snapshot\" \\")
    print("      --item  \"Wafer pricing;product-by-product\" \\")
    print("      --item  \"Outlook;next-quarter view\"")
    print()
    print("  --- add-section -----------------------------------------------")
    print("  Section divider page. --number defaults to 01,02,... auto.")
    print("    pptx_helper.py add-section FILE.pptx --title \"Market Overview\"")
    print("    pptx_helper.py add-section FILE.pptx --title \"Pricing\" --number \"02\" --intro \"Detail by SKU\"")
    print()
    print("  --- add-text --------------------------------------------------")
    print("  Plain bullet content (1..5 bullets; --bullet repeated)")
    print("    pptx_helper.py add-text FILE.pptx \\")
    print("      --title  \"Key Findings\" \\")
    print("      --bullet \"NAND prices flat WoW\" \\")
    print("      --bullet \"DRAM contract +1.2%\" \\")
    print("      --bullet \"Demand still soft into Q2\"")
    print()
    print("  --- add-table -------------------------------------------------")
    print("  Full-width data table. Headers and rows are comma-separated.")
    print("  Max 5 columns x 8 rows; longer auto-truncated with [WARN].")
    print("    pptx_helper.py add-table FILE.pptx \\")
    print("      --title   \"NAND Wafer\" \\")
    print("      --headers \"Product,Spot,WoW\" \\")
    print("      --row     \"1Tb QLC,$27.00,0%\" \\")
    print("      --row     \"1Tb TLC,$29.00,+0.5%\" \\")
    print("      --row     \"512Gb TLC,$15.20,-0.2%\"")
    print()
    print("  --- add-stats -------------------------------------------------")
    print("  2..4 big-number stat boxes. --stat \"value;label[;desc]\"")
    print("    pptx_helper.py add-stats FILE.pptx \\")
    print("      --title \"Snapshot\" \\")
    print("      --stat  \"$27;QLC 1Tb;spot\" \\")
    print("      --stat  \"0%;WoW change\" \\")
    print("      --stat  \"3;Active SKUs\"")
    print()
    print("  --- add-compare -----------------------------------------------")
    print("  Two-column A vs B card. Each side has a title + 2..5 points.")
    print("    pptx_helper.py add-compare FILE.pptx \\")
    print("      --title       \"QLC vs TLC\" \\")
    print("      --left-title  \"QLC\" \\")
    print("      --left-point  \"Higher density\" \\")
    print("      --left-point  \"Lower endurance\" \\")
    print("      --right-title \"TLC\" \\")
    print("      --right-point \"Mainstream\" \\")
    print("      --right-point \"Better endurance\"")
    print()
    print("  --- add-timeline ----------------------------------------------")
    print("  3..5 step horizontal timeline. --step \"label;desc\"")
    print("    pptx_helper.py add-timeline FILE.pptx \\")
    print("      --title \"Roadmap\" \\")
    print("      --step  \"Q1;Plan\" \\")
    print("      --step  \"Q2;Build\" \\")
    print("      --step  \"Q3;Ship\"")
    print()
    print("  --- add-text-table --------------------------------------------")
    print("  Mixed: bullets on one side + table on the other.")
    print("  --pos = right (default) | left | top | bottom")
    print("    pptx_helper.py add-text-table FILE.pptx \\")
    print("      --title   \"Q1 Wafer Trend\" --pos right \\")
    print("      --bullet  \"QLC flat\" --bullet \"TLC slight up\" \\")
    print("      --headers \"Product,Price\" \\")
    print("      --row     \"QLC,$27\" --row \"TLC,$29\"")
    print()
    print("  --- add-text-image --------------------------------------------")
    print("  Mixed: bullets + image. --pos = right (default) | left.")
    print("  If --image PATH is omitted or missing, draws a placeholder box.")
    print("    pptx_helper.py add-text-image FILE.pptx \\")
    print("      --title \"Test rig\" --pos right \\")
    print("      --bullet \"PCIe 5.0\" --bullet \"M.2 form factor\" \\")
    print("      --image  /sdcard/photo.png --caption \"Lab setup\"")
    print()
    print("  --- add-summary -----------------------------------------------")
    print("  Closing page. 2..5 takeaways + optional contact line.")
    print("    pptx_helper.py add-summary FILE.pptx \\")
    print("      --title    \"Conclusion\" \\")
    print("      --takeaway \"Prices stable\" \\")
    print("      --takeaway \"Demand soft\" \\")
    print("      --takeaway \"Watch Q2\" \\")
    print("      --contact  \"research@example.com\"")
    print()
    print("=== STEP 4: edit text in place (optional) ===")
    print("  pptx_helper.py inspect FILE.pptx               # see (slide,shape) addresses")
    print("  # set-text FILE SLIDE SHAPE TEXT       (positional)")
    print("  pptx_helper.py set-text FILE.pptx 1 3 \"New title\"")
    print("  # set-cell FILE SLIDE SHAPE ROW COL TEXT  (positional)")
    print("  pptx_helper.py set-cell FILE.pptx 2 8 1 2 \"$32.00\"")
    print()
    print("=== HARD LIMITS (auto-clipped, never fail) ===")
    print(f"  title <= {_LIM['title']} chars      bullet <= {_LIM['bullet']} chars      table <= "
          f"{_LIM['headers_max']}x{_LIM['rows_max']} cells")
    print(f"  bullets list 1..{_LIM['bullets_max']}        stats list "
          f"{_LIM['stats_min']}..{_LIM['stats_max']}        steps list "
          f"{_LIM['steps_min']}..{_LIM['steps_max']}")
    print(f"  toc items {_LIM['toc_items_min']}..{_LIM['toc_items_max']}        takeaways "
          f"{_LIM['takeaways_min']}..{_LIM['takeaways_max']}        compare points "
          f"{_LIM['compare_points_min']}..{_LIM['compare_points_max']}")
    print()
    _next_step([
        "Build the deck by chaining commands. Minimum viable deck = 4 slides:",
        "  pptx_helper.py new-deck    ${WORKSPACE}/deck.pptx --theme business",
        "  pptx_helper.py add-cover   ${WORKSPACE}/deck.pptx --title \"...\" --subtitle \"...\"",
        "  pptx_helper.py add-table   ${WORKSPACE}/deck.pptx --title \"...\" --headers \"...\" --row \"...\"",
        "  pptx_helper.py add-summary ${WORKSPACE}/deck.pptx --title \"...\" --takeaway \"...\" --takeaway \"...\"",
        "All inputs are auto-clipped to fit the slide; no plan.json, no dry-run.",
    ])


def cmd_plan(args):
    plan_path = args.plan_pos or args.plan_opt
    if not plan_path:
        raise ValueError("plan path is required (positional or --plan PATH)")
    plan = _load_json(plan_path)
    issues = _validate_plan(plan)
    if issues:
        print("[VALIDATION] issues found:")
        for it in issues:
            print(f"  - {it}")
    else:
        print("[VALIDATION] OK")
    if args.dry_run:
        print("[DRY-RUN] no file written.")
        if issues:
            _next_step([
                "Fix the issues above in plan.json (slot names/types/counts), then re-run:",
                f"  python pptx_helper.py plan {plan_path} --dry-run",
                "Run `pptx_helper.py catalog` to see the exact slot schema.",
            ])
            sys.exit(1)
        _next_step([
            "Plan passed validation. Build with:",
            f"  python pptx_helper.py plan {plan_path} --out ${{WORKSPACE}}/deck.pptx",
        ])
        return
    if issues:
        # Hard reject build on any structural error. Soft warning (file-not-found
        # for image_path) is the only kind allowed to pass through.
        hard = [i for i in issues if "will render placeholder" not in i]
        if hard:
            # Loud, model-friendly failure marker. Engines that swallow exit
            # codes (e.g. embedded Python runtimes) will still see this banner
            # in stdout and can pattern-match on it. This complements sys.exit(1).
            print("[BUILD_FAILED] plan has validation issues; "
                  "no .pptx file was written.")
            print("[ERROR] plan has validation issues; not building. "
                  "Use --dry-run to debug, then re-run with --out.")
            sys.exit(1)
    if not args.out:
        raise ValueError("--out PATH is required unless --dry-run is set")
    theme, n = _build_from_plan(plan, args.out)
    print(f"[OK] plan -> {args.out} (theme={theme}, slides={n})")
    _next_step([
        f"Deck generated at {args.out}",
        f"Inspect: python pptx_helper.py inspect {args.out}",
        "Tweak text in place: python pptx_helper.py set-text ... --slide N --shape K --text \"...\"",
    ])


class SmartArgumentParser(argparse.ArgumentParser):
    def error(self, message):
        sys.stderr.write(f"[ERROR] {message}\n")
        # Print SKILL.md so the model can pick a different subcommand or fix
        # the args of the same one without re-reading any other manual.
        _print_commands_menu(stream=sys.stderr)
        sys.exit(2)


# Hallucinated subcommand -> intent-aware "did you mean" hint. Argparse's
# default invalid-choice error already lists every legal cmd, but the model
# still picks the wrong one when the name does not match its intuition (e.g.
# new-ppt, create, new-pptx). Map each known hallucination to the canonical
# replacement so the model sees a directed fix instead of having to guess.
_PPTX_DID_YOU_MEAN = {
    "new-ppt":    "new-deck FILE.pptx --theme business [--template PATH]",
    "new-pptx":   "new-deck FILE.pptx --theme business [--template PATH]",
    "new":        "new-deck FILE.pptx --theme business",
    "create":     "new-deck FILE.pptx --theme business",
    "create-deck": "new-deck FILE.pptx --theme business",
    "add-slide":  "add-text / add-table / add-stats / add-cover / ... (one slide per add-*)",
    "set-slide":  "set-text FILE SLIDE SHAPE \"new text\" (in-place) "
                  "or inspect FILE first to find the (slide, shape) target",
}


def _intercept_hallucinated_cmd(known_cmds):
    """If sys.argv[1] is a well-known wrong name, print a directed fix and
    exit(2) BEFORE argparse runs. Falls through silently otherwise."""
    if len(sys.argv) < 2:
        return
    cmd = sys.argv[1]
    if cmd in known_cmds or cmd.startswith("-"):
        return
    if cmd in _PPTX_DID_YOU_MEAN:
        sys.stderr.write(
            f"[ERROR] unknown subcommand '{cmd}'. "
            f"Did you mean:  {_PPTX_DID_YOU_MEAN[cmd]}\n"
        )
        _print_commands_menu(stream=sys.stderr)
        sys.exit(2)


# ===========================================================================
# Argparse wiring
# ===========================================================================
_PPTX_KNOWN_CMDS = {
    "inspect", "extract", "set-text", "set-cell",
    "catalog", "new-deck",
    "add-cover", "add-toc", "add-section", "add-text", "add-table",
    "add-stats", "add-compare", "add-timeline",
    "add-text-table", "add-text-image", "add-summary",
}


def main():
    # Intercept well-known wrong subcommand names BEFORE argparse so the model
    # gets a directed fix (e.g. new-ppt -> new-deck) instead of a generic
    # "invalid choice" listing.
    _intercept_hallucinated_cmd(_PPTX_KNOWN_CMDS)

    ap = SmartArgumentParser(prog="pptx_helper",
                             description="Programmatic helper for .pptx files.")
    sp = ap.add_subparsers(dest="cmd", required=True)

    p_ins = sp.add_parser("inspect", help="Per-slide shape listing for LLM")
    p_ins.add_argument("file")
    p_ins.set_defaults(func=cmd_inspect)

    p_ext = sp.add_parser("extract", help="Plain-text dump")
    p_ext.add_argument("file")
    p_ext.set_defaults(func=cmd_extract)

    # set-text FILE SLIDE SHAPE TEXT  (canonical positional form)
    p_st = sp.add_parser("set-text",
                         help="Replace shape text: FILE SLIDE SHAPE TEXT")
    p_st.add_argument("file")
    p_st.add_argument("slide_pos", nargs="?", default=None,
                      help="slide index (0-based) from `inspect`")
    p_st.add_argument("shape_pos", nargs="?", default=None,
                      help="shape index (0-based) from `inspect`")
    p_st.add_argument("text_pos", nargs="?", default=None,
                      help="new text; '\\n' for line breaks")
    # Legacy flags (hidden; still functional).
    p_st.add_argument("--slide", type=int, default=None, help=argparse.SUPPRESS)
    p_st.add_argument("--shape", type=int, default=None, help=argparse.SUPPRESS)
    p_st.add_argument("--text", default=None, help=argparse.SUPPRESS)
    p_st.add_argument("--out", "--output", dest="out", default=None,
                      help=argparse.SUPPRESS)
    p_st.set_defaults(func=cmd_set_text)

    # set-cell FILE SLIDE SHAPE ROW COL TEXT  (canonical positional form)
    p_sc = sp.add_parser("set-cell",
                         help="Replace table cell text: FILE SLIDE SHAPE ROW COL TEXT")
    p_sc.add_argument("file")
    p_sc.add_argument("slide_pos", nargs="?", default=None,
                      help="slide index (0-based)")
    p_sc.add_argument("shape_pos", nargs="?", default=None,
                      help="shape index (0-based)")
    p_sc.add_argument("row_pos", nargs="?", default=None,
                      help="row index (0-based)")
    p_sc.add_argument("col_pos", nargs="?", default=None,
                      help="column index (0-based)")
    p_sc.add_argument("text_pos", nargs="?", default=None,
                      help="new cell text")
    # Legacy flags (hidden; still functional).
    p_sc.add_argument("--slide", type=int, default=None, help=argparse.SUPPRESS)
    p_sc.add_argument("--shape", type=int, default=None, help=argparse.SUPPRESS)
    p_sc.add_argument("--row", type=int, default=None, help=argparse.SUPPRESS)
    p_sc.add_argument("--col", type=int, default=None, help=argparse.SUPPRESS)
    p_sc.add_argument("--text", default=None, help=argparse.SUPPRESS)
    p_sc.add_argument("--out", "--output", dest="out", default=None,
                      help=argparse.SUPPRESS)
    p_sc.set_defaults(func=cmd_set_cell)

    # ------------------------------------------------------------------
    # Direct CLI deck builder (RECOMMENDED). Each add-* command appends
    # one slide to FILE.pptx. No plan.json, no schema, no dry-run.
    # ------------------------------------------------------------------
    p_cat = sp.add_parser("catalog",
                          help="Cookbook of themes + add-* commands (run this first)")
    p_cat.set_defaults(func=cmd_catalog)

    p_nd = sp.add_parser("new-deck",
                         help=("Init a deck FILE.pptx; optionally inherit a "
                               "user template's master/cover via --template"))
    p_nd.add_argument("file", help="output .pptx path")
    p_nd.add_argument("--theme", default="business",
                      help=f"color palette: {list(PALETTES)}")
    p_nd.add_argument("--template", default=None,
                      help=("optional .pptx/.potx path; its slide master "
                            "(logo, brand fonts, footer) is inherited and "
                            "any slides it already contains are kept; our "
                            "add-* commands append AFTER them"))
    p_nd.set_defaults(func=cmd_new_deck)

    p_ac = sp.add_parser("add-cover", help="Append a cover slide")
    p_ac.add_argument("file")
    p_ac.add_argument("--title", required=True)
    p_ac.add_argument("--subtitle", default="")
    p_ac.add_argument("--meta", default="")
    p_ac.set_defaults(func=cmd_add_cover)

    p_at = sp.add_parser("add-toc", help="Append a table-of-contents slide")
    p_at.add_argument("file")
    p_at.add_argument("--title", default="Table of Contents")
    p_at.add_argument("--item", action="append", default=[],
                      help="repeatable; format \"Title\" or \"Title;Description\"")
    p_at.set_defaults(func=cmd_add_toc)

    p_as = sp.add_parser("add-section", help="Append a section divider slide")
    p_as.add_argument("file")
    p_as.add_argument("--title", required=True)
    p_as.add_argument("--number", default=None,
                      help="optional, e.g. \"01\". Auto-generated if omitted.")
    p_as.add_argument("--intro", default="")
    p_as.set_defaults(func=cmd_add_section)

    p_atx = sp.add_parser("add-text", help="Append a plain bullet content slide")
    p_atx.add_argument("file")
    p_atx.add_argument("--title", required=True)
    p_atx.add_argument("--bullet", action="append", default=[],
                       help=f"repeatable; 1..{_LIM['bullets_max']} bullets recommended")
    p_atx.set_defaults(func=cmd_add_text)

    p_atb = sp.add_parser("add-table", help="Append a full-width table slide")
    p_atb.add_argument("file")
    # Positional fallback aligns with xlsx add-sheet style:
    #   add-table FILE TITLE HEADERS ROW1 ROW2 ...
    p_atb.add_argument("pos_args", nargs="*",
                       help="positional fallback: TITLE HEADERS [ROW1 ROW2 ...]")
    p_atb.add_argument("--title", default=None)
    p_atb.add_argument("--headers", default=None,
                       help="comma-separated header cells (max %d)" % _LIM["headers_max"])
    p_atb.add_argument("--row", action="append", default=[],
                       help="repeatable; comma-separated cells (max %d rows)" % _LIM["rows_max"])
    p_atb.set_defaults(func=cmd_add_table)

    p_ast = sp.add_parser("add-stats", help="Append a 2..4 big-number stats slide")
    p_ast.add_argument("file")
    p_ast.add_argument("--title", required=True)
    p_ast.add_argument("--stat", action="append", default=[],
                       help="repeatable; format \"value;label\" or \"value;label;desc\"")
    p_ast.set_defaults(func=cmd_add_stats)

    p_acm = sp.add_parser("add-compare", help="Append a 2-column compare slide")
    p_acm.add_argument("file")
    p_acm.add_argument("--title", required=True)
    p_acm.add_argument("--left-title", default="A")
    p_acm.add_argument("--left-point", action="append", default=[],
                       help="repeatable; bullet points in the LEFT column")
    p_acm.add_argument("--right-title", default="B")
    p_acm.add_argument("--right-point", action="append", default=[],
                       help="repeatable; bullet points in the RIGHT column")
    p_acm.set_defaults(func=cmd_add_compare)

    p_atl = sp.add_parser("add-timeline", help="Append a 3..5 step timeline slide")
    p_atl.add_argument("file")
    p_atl.add_argument("--title", required=True)
    p_atl.add_argument("--step", action="append", default=[],
                       help="repeatable; format \"label;desc\"")
    p_atl.set_defaults(func=cmd_add_timeline)

    p_att = sp.add_parser("add-text-table",
                          help="Append a mixed bullets+table slide")
    p_att.add_argument("file")
    p_att.add_argument("--title", required=True)
    p_att.add_argument("--pos", default="right",
                       help="table position: right (default) | left | top | bottom")
    p_att.add_argument("--bullet", action="append", default=[],
                       help="repeatable; bullets in the text panel")
    p_att.add_argument("--headers", required=True,
                       help="comma-separated header cells")
    p_att.add_argument("--row", action="append", default=[],
                       help="repeatable; comma-separated cells")
    p_att.set_defaults(func=cmd_add_text_table)

    p_ati = sp.add_parser("add-text-image",
                          help="Append a mixed bullets+image slide")
    p_ati.add_argument("file")
    p_ati.add_argument("--title", required=True)
    p_ati.add_argument("--pos", default="right",
                       help="image position: right (default) | left")
    p_ati.add_argument("--bullet", action="append", default=[],
                       help="repeatable; bullets in the text panel")
    p_ati.add_argument("--image", default="",
                       help="image file path; placeholder drawn if missing")
    p_ati.add_argument("--caption", default="")
    p_ati.set_defaults(func=cmd_add_text_image)

    p_asm = sp.add_parser("add-summary", help="Append a closing takeaways slide")
    p_asm.add_argument("file")
    p_asm.add_argument("--title", default="Key Takeaways")
    p_asm.add_argument("--takeaway", action="append", default=[],
                       help=f"repeatable; {_LIM['takeaways_min']}..{_LIM['takeaways_max']} items")
    p_asm.add_argument("--contact", default="")
    p_asm.set_defaults(func=cmd_add_summary)

    args = ap.parse_args()
    try:
        args.func(args)
        # SKILL.md tail: the model reads this on every reply to know what
        # commands exist and what the next step could be. We do NOT presume
        # to suggest a specific next step -- that is the task-level reasoner's
        # call, not the helper's.
        _print_commands_menu()
    except Exception as e:
        sys.stderr.write(f"[ERROR] {type(e).__name__}: {e}\n")
        _print_commands_menu(stream=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
